#!/usr/bin/env python3
"""Test suite for check_pptx.py — the PPTX QA gate.

Runs standalone (no pytest needed):  python3 scripts/test_check_pptx.py
Also pytest-compatible:               pytest scripts/test_check_pptx.py

Builds two decks in-memory with python-pptx (pptxgenjs is not required — these are
fixtures for the *checker*, not the kit): a clean deck that should score >= 80, and a
deliberately broken one exercising every finding category, which must score < 80.
"""
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)

import check_pptx  # noqa: E402

ASSETS = os.path.join(HERE, "..", "assets")
LOGO = os.path.join(ASSETS, "logo", "aws_logo.png")


def _add_footer(slide, page_num):
    tb = slide.shapes.add_textbox(Inches(0.92), Inches(7.06), Inches(8.2), Inches(0.3))
    tb.text_frame.text = "© 2026, Amazon Web Services, Inc. or its affiliates. All rights reserved."
    for r in tb.text_frame.paragraphs[0].runs:
        r.font.size = Pt(8)
        r.font.name = "Pretendard"
    pn = slide.shapes.add_textbox(Inches(12.5), Inches(7.06), Inches(0.4), Inches(0.3))
    pn.text_frame.text = str(page_num)
    for r in pn.text_frame.paragraphs[0].runs:
        r.font.size = Pt(9)
        r.font.name = "Pretendard"


def _add_text(slide, x, y, w, h, text, size=14, font="Pretendard"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    tb.text_frame.text = text
    for para in tb.text_frame.paragraphs:
        if not para.runs:
            # empty text still gets one implicit run once we set .text below
            continue
        for r in para.runs:
            r.font.size = Pt(size)
            r.font.name = font
    return tb


def _blank_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _good_deck():
    prs = _blank_prs()
    layout = prs.slide_layouts[6]  # blank
    for i in range(1, 4):
        s = prs.slides.add_slide(layout)
        _add_text(s, 0.92, 0.6, 11, 0.8, f"슬라이드 {i} 제목", size=30)
        _add_text(s, 0.92, 1.6, 11, 1.0, "짧은 본문 텍스트입니다.", size=14)
        _add_footer(s, i + 1)
    return prs


def _bad_deck():
    prs = _blank_prs()
    layout = prs.slide_layouts[6]

    # Slide 1: text overflow (20 lines of Korean text crammed into a 1" box) + Arial run + TODO
    s1 = prs.slides.add_slide(layout)
    long_text = "\n".join(f"이것은 넘치는 텍스트 라인 번호 {n} 입니다" for n in range(20))
    box = _add_text(s1, 1.0, 1.0, 2.0, 1.0, long_text, size=18, font="Arial")
    for para in box.text_frame.paragraphs:
        for r in para.runs:
            r.font.name = "Arial"
    todo_box = _add_text(s1, 1.0, 5.0, 3.0, 0.4, "TODO: fill this in", size=12)
    # no footer on slide 1 -> missing-footer finding too

    # Slide 2: two overlapping pictures + duplicate page number (2 after slide1's implicit none)
    s2 = prs.slides.add_slide(layout)
    s2.shapes.add_picture(LOGO, Inches(2.0), Inches(2.0), width=Inches(2.0), height=Inches(1.2))
    s2.shapes.add_picture(LOGO, Inches(2.5), Inches(2.3), width=Inches(2.0), height=Inches(1.2))
    _add_footer(s2, 2)

    # Slide 3: duplicate page number (2 again -> regression/duplicate) + off-canvas shape
    s3 = prs.slides.add_slide(layout)
    _add_text(s3, 13.0, 0.2, 2.0, 0.5, "off canvas", size=12)
    _add_footer(s3, 2)

    return prs


def test_good_deck_passes():
    r = check_pptx.analyze(_good_deck())
    assert check_pptx.gate_pass(r, 80), f"good deck scored {r['score']}: {r['findings']}"


def test_bad_deck_fails():
    r = check_pptx.analyze(_bad_deck())
    assert not check_pptx.gate_pass(r, 80), f"bad deck passed the gate: {r['findings']}"
    joined = " ".join(r["findings"])
    assert "overflow" in joined, r["findings"]
    assert "overlapping" in joined, r["findings"]
    assert "footer" in joined, r["findings"]
    assert "page-number" in joined, r["findings"]
    assert "Pretendard" in joined, r["findings"]
    assert "placeholder" in joined, r["findings"]


def test_offcanvas_detected():
    r = check_pptx.analyze(_bad_deck())
    assert any("outside the slide canvas" in f for f in r["findings"]), r["findings"]


def test_geometry_finding_fails_gate_despite_high_score():
    # one small off-canvas shape -> score stays high (>=80) but the gate must still fail
    prs = _blank_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, 0.92, 0.6, 11, 0.8, "제목", size=30)
    _add_footer(s, 2)
    _add_text(s, 13.2, 0.2, 1.5, 0.4, "off", size=12)  # off-canvas -> -5 geometry
    r = check_pptx.analyze(prs)
    assert r["score"] >= 80, f"expected high score, got {r['score']}"
    assert r["geometry_findings"] >= 1
    assert not check_pptx.gate_pass(r, 80), "geometry finding must fail the gate regardless of score"


def test_body_mention_does_not_satisfy_footer():
    # 'Amazon Web Services' in body text (not the footer band) must NOT count as a footer
    prs = _blank_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, 0.92, 1.6, 11, 1.0, "Amazon Web Services는 클라우드입니다.", size=14)
    r = check_pptx.analyze(prs)
    assert any("footer" in f for f in r["findings"]), r["findings"]


def test_no_page_numbers_flagged():
    prs = _blank_prs()
    for _ in range(2):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _add_text(s, 0.92, 1.0, 11, 3.0, "본문", size=14)
        # footer copyright but NO page number
        tb = s.shapes.add_textbox(Inches(0.92), Inches(7.06), Inches(8), Inches(0.3))
        tb.text_frame.text = "© 2026, Amazon Web Services, Inc. All rights reserved."
    r = check_pptx.analyze(prs)
    assert any("no page number" in f for f in r["findings"]), r["findings"]


def test_caption_on_icon_overlap_detected():
    # a text box sitting on top of a picture (the defect a too-short cell height makes)
    # must be caught even though the two shapes are different kinds
    prs = _blank_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_picture(LOGO, Inches(3.0), Inches(2.0), width=Inches(0.62), height=Inches(0.62))
    _add_text(s, 3.0, 2.0, 1.5, 0.5, "겹치는 캡션", size=11)  # same top-left as the icon
    _add_footer(s, 2)
    r = check_pptx.analyze(prs)
    assert any("overlapping" in f for f in r["findings"]), r["findings"]


def test_step_marker_on_icon_not_flagged():
    # arch_kit's stepMarker is a 0.34" numbered badge intentionally overlaying an icon
    # corner — badge-sized text boxes must be exempt from picture x text overlap
    prs = _blank_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_picture(LOGO, Inches(3.0), Inches(2.0), width=Inches(0.62), height=Inches(0.62))
    _add_text(s, 3.14, 1.9, 0.34, 0.34, "5", size=12)  # marker on the icon's top-right
    _add_footer(s, 2)
    r = check_pptx.analyze(prs)
    assert not any("overlapping" in f for f in r["findings"]), r["findings"]


def test_soft_break_counts_as_line():
    # "\n" in a small box should push the wrapped height over the box -> overflow
    prs = _blank_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, 1.0, 1.0, 3.0, 0.4, "Elastic Load\nBalancing\n세 번째 줄\n네 번째 줄", size=18)
    _add_footer(s, 2)
    r = check_pptx.analyze(prs)
    assert any("overflow" in f for f in r["findings"]), r["findings"]


def test_large_image_fully_offcanvas_flagged():
    # a big picture (>80% of slide area) placed entirely off-canvas must still flag —
    # the background exemption is coverage-based, not area-based
    prs = _blank_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_footer(s, 2)
    s.shapes.add_picture(LOGO, Inches(14.0), Inches(1.0), width=Inches(12.0), height=Inches(6.0))
    r = check_pptx.analyze(prs)
    assert any("outside the slide canvas" in f for f in r["findings"]), r["findings"]


def test_full_bleed_background_exempt():
    # a picture that actually covers the canvas (bleeding slightly past edges) is exempt
    prs = _blank_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_footer(s, 2)
    s.shapes.add_picture(LOGO, Inches(-0.2), Inches(-0.2), width=Inches(13.7), height=Inches(7.9))
    r = check_pptx.analyze(prs)
    assert not any("outside the slide canvas" in f for f in r["findings"]), r["findings"]


def test_group_shapes_fail_gate():
    # native GROUP children can't be inspected, so a grouped deck must fail the gate
    # (forces manual review) rather than silently passing with geometry_findings=0
    prs = _blank_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_footer(s, 2)
    gs = s.shapes.add_group_shape()
    gs.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5)).text_frame.text = "in group"
    r = check_pptx.analyze(prs)
    assert r["geometry_findings"] >= 1, r["findings"]
    assert not check_pptx.gate_pass(r, 80), "grouped deck must fail the gate"
    assert any("GROUP" in f for f in r["findings"]), r["findings"]


def test_table_cell_placeholder_detected():
    # placeholder text inside a table cell must be caught (tables were invisible before)
    prs = _blank_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_footer(s, 2)
    gf = s.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2))
    gf.table.cell(0, 0).text = "TODO: fill this cell"
    r = check_pptx.analyze(prs)
    assert any("placeholder" in f for f in r["findings"]), r["findings"]


def test_table_cell_overflow_detected():
    # a tiny table cell crammed with text must flag as overflow (geometry finding)
    prs = _blank_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_footer(s, 2)
    gf = s.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(1.2), Inches(0.3))
    gf.table.cell(0, 0).text = "\n".join(f"line {i}" for i in range(12))
    r = check_pptx.analyze(prs)
    assert any("overflow" in f for f in r["findings"]), r["findings"]


def test_merged_header_cell_no_false_overflow():
    # a header merged across 3 columns must be measured at its full spanned width —
    # naive per-position iteration re-checked it at single-column size (false overflow)
    prs = _blank_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_footer(s, 2)
    gf = s.shapes.add_table(2, 3, Inches(1), Inches(1), Inches(9), Inches(1.6))
    origin = gf.table.cell(0, 0)
    origin.merge(gf.table.cell(0, 2))
    origin.text = "가속 컴퓨팅 포트폴리오 — GPU 및 AWS ML 가속기 헤더"  # fits 9", not 3"
    r = check_pptx.analyze(prs)
    assert not any("overflow" in f for f in r["findings"]), r["findings"]


def test_circled_digit_footer_no_crash():
    prs = _blank_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(s, 0.92, 1.0, 11, 3.0, "본문", size=14)
    tb = s.shapes.add_textbox(Inches(12.5), Inches(7.06), Inches(0.4), Inches(0.3))
    tb.text_frame.text = "\u2461"  # circled digit two: isdigit() True, int() would crash
    r = check_pptx.analyze(prs)  # must not raise
    assert isinstance(r["score"], int), r


def test_table_overlap_detected():
    prs = _blank_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_footer(s, 2)
    s.shapes.add_table(2, 2, Inches(2), Inches(2), Inches(4), Inches(2))
    s.shapes.add_picture(LOGO, Inches(3), Inches(2.5), width=Inches(2), height=Inches(1.2))
    r = check_pptx.analyze(prs)
    assert any("overlapping" in f for f in r["findings"]), r["findings"]


def test_empty_deck_fails_gate():
    prs = _blank_prs()  # zero slides
    r = check_pptx.analyze(prs)
    assert not check_pptx.gate_pass(r, 80), "empty deck must fail the gate"
    assert any("no slides" in f for f in r["findings"]), r["findings"]


def test_threshold_noninteger_exits_2():
    import subprocess
    p = subprocess.run([sys.executable, os.path.join(HERE, "check_pptx.py"),
                        os.path.join(ASSETS, "logo", "aws_logo.png"), "--threshold", "abc"],
                       capture_output=True, text=True)
    assert p.returncode == 2, f"expected exit 2 on bad threshold, got {p.returncode}: {p.stdout}{p.stderr}"


# --- standalone runner (mirrors test_layout_aws.py) ---
def _run():
    tests = sorted((n, f) for n, f in globals().items()
                    if n.startswith("test_") and callable(f))
    passed = failed = 0
    for name, fn in tests:
        try:
            fn()
            print(f"  ✅ {name}")
            passed += 1
        except Exception as e:  # noqa: BLE001
            print(f"  ❌ {name}: {type(e).__name__}: {e}")
            failed += 1
    print(f"\n{passed} passed, {failed} failed ({len(tests)} total)")
    return 1 if failed else 0


if __name__ == "__main__":
    sys.exit(_run())
