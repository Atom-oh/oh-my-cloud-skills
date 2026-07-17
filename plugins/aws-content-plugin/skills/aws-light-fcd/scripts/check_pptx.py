#!/usr/bin/env python3
"""Programmatic QA for a .pptx deck — the "is it clean enough to ship?" gate.

Mirrors architecture-diagram's `lint_layout.py`: this is the deterministic check that
catches what an LLM builder gets subtly wrong and a human would fix on sight. Two
families:

  GEOMETRY (does content fit the canvas?)
  - text overflow      (estimated wrapped-text height exceeds the shape's box)
  - element overlap     (two picture boxes, or two text boxes, intersect)
  - off-canvas          (a shape's box falls outside the slide bounds)

  DESIGN (does it look finished, like a hand-crafted deck?)
  - footer presence      (every slide should carry the copyright line)
  - page-number sanity   (footer page numbers must be strictly increasing, no dupes)
  - font consistency      (every run should use the deck's one font — Pretendard)
  - minimum font size     (nothing smaller than 8pt — the footer's own floor)
  - placeholder text      (TODO/TBD/lorem/??? left behind by a half-finished build)

It prints a deck SCORE (0-100, with a geometry/design breakdown) and exits non-zero
when the gate fails, so it can be used as a hard pre-delivery gate — the PPTX
equivalent of the "convert to PDF and eyeball it" step, but automatable and exact.

THE GATE: pass = score >= threshold AND zero GEOMETRY findings. Geometry findings are
never scored away — a 95-point deck with one overflowing text box still fails, because
content-review-agent treats any geometry finding as Critical (same rule, one place).
Design findings only lower the score.

Trust boundary: this parses untrusted OOXML via python-pptx/lxml — treat decks from
outside the kit as untrusted input and run it in your normal sandbox, not with
elevated privileges.

Usage:
    python3 check_pptx.py <file.pptx>                 # score + findings
    python3 check_pptx.py <file.pptx> --threshold 85   # custom gate (default 80)
    python3 check_pptx.py <file.pptx> --json           # machine-readable

Exit 0 = passes the gate. 1 = gate failed (fix before delivery). 2 = usage/read error.
"""
import re
import sys
import json as _json

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# --- Tunables ---
DEFAULT_THRESHOLD = 80
EMU_PER_IN = 914400
DEFAULT_MARGIN_LR_EMU = 91440    # OOXML bodyPr default lIns/rIns (0.1")
DEFAULT_MARGIN_TB_EMU = 45720    # OOXML bodyPr default tIns/bIns (0.05")
DEFAULT_FONT_PT = 12             # fallback when no run/paragraph size is set
LINE_HEIGHT_MULT = 1.25          # matches this kit's typical lineSpacingMultiple
OVERFLOW_TOL = 1.10              # 10% slack before an overflow counts
OVERLAP_AREA_RATIO = 0.15        # intersection > 15% of the smaller shape -> flag
BG_AREA_RATIO = 0.80             # a picture covering >80% of the slide is a background, not content
BADGE_MAX_IN = 0.40              # text box <= this in both dims = badge/marker (may overlay a picture)
CANVAS_TOL_IN = 0.02
FONT_NAME = "Pretendard"
COPYRIGHT_MARK = "amazon web services"
FOOTER_BAND_IN = 6.5             # y-position past which text is "in the footer row"
PLACEHOLDER_RE = re.compile(r"\bTODO\b|\bTBD\b|lorem ipsum|\bXXX\b|placeholder|\?{3}", re.I)


def _iter_shapes(shapes):
    """Flatten a shape tree. NOTE: python-pptx group-child coordinates live in the
    group's own child-coordinate space, not the slide's — correctly re-projecting them
    needs the group's transform/scale math. This kit (deck_kit.js/arch_kit.js) never
    emits native PPTX groups (it composes with plain addShape/addText/addImage calls),
    so we don't do that projection; a GROUP shape is checked as its own bbox only and
    its children are skipped by ALL checks (geometry, font, footer, placeholder).
    Because that would let a grouped external deck pass with nothing inspected, analyze()
    turns GROUP presence into a blocking geometry finding (fails the gate, forces a manual
    review) rather than a silent pass."""
    for sh in shapes:
        yield sh
        # GROUP children intentionally not recursed into (see docstring) — the group's
        # own bbox was just yielded and analyze() emits a per-deck note when groups exist.


def _bbox_in(sh):
    if sh.left is None or sh.top is None or sh.width is None or sh.height is None:
        return None
    return (sh.left / EMU_PER_IN, sh.top / EMU_PER_IN, sh.width / EMU_PER_IN, sh.height / EMU_PER_IN)


def _has_text(sh):
    return bool(getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip())


def _kind(sh, slide_area_sqin):
    try:
        st = sh.shape_type
    except (ValueError, NotImplementedError):
        st = None
    if st == MSO_SHAPE_TYPE.PICTURE:
        box = _bbox_in(sh)
        if box and (box[2] * box[3]) / slide_area_sqin > BG_AREA_RATIO:
            return None  # full-bleed background image, not "content" for overlap purposes
        return "picture"
    if _has_text(sh):
        return "text"
    return None


def _wrap_lines(text, fs_pt, budget_pt):
    """Estimated line count. Hard breaks first (python-pptx returns '\\v' for <a:br>,
    which is what PptxGenJS emits for a '\\n' in text), then character-by-character wrap
    of each segment. Deliberate simplification: word-wrap (not char-wrap) would be more
    accurate for space-delimited English, but this deck is Korean-first (no spaces to
    break on) — char-wrap is the simpler estimator that works for both. Upgrade to real
    word-wrap if false positives show up on English-heavy decks."""
    if not text:
        return 1
    total = 0
    for seg in re.split(r"[\n\v]", text):
        total += 1  # each hard-break segment is at least one line
        cur = 0.0
        for ch in seg:
            # CJK glyphs are ~1em wide; Latin/digits/punctuation average ~0.55em (same
            # factor used by deck_kit.js's estTextW, so the two stay in sync).
            cw = fs_pt * (1.0 if ord(ch) > 0x2E80 else 0.55)
            if cur + cw > budget_pt and cur > 0:
                total += 1
                cur = cw
            else:
                cur += cw
    return total


def _para_font_pt(para):
    # largest run in the paragraph — mixed-size runs overflow at the biggest glyph
    sizes = [r.font.size.pt for r in para.runs if r.font.size]
    if sizes:
        return max(sizes)
    if para.font.size:
        return para.font.size.pt
    return DEFAULT_FONT_PT


def _check_overflow(shapes):
    n, example = 0, None
    for sh in shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        tf = sh.text_frame
        if not tf.text.strip():
            continue
        box = _bbox_in(sh)
        if not box or box[2] <= 0 or box[3] <= 0:
            continue
        _, _, w_in, h_in = box
        ml = (tf.margin_left if tf.margin_left is not None else DEFAULT_MARGIN_LR_EMU) / EMU_PER_IN
        mr = (tf.margin_right if tf.margin_right is not None else DEFAULT_MARGIN_LR_EMU) / EMU_PER_IN
        mt = (tf.margin_top if tf.margin_top is not None else DEFAULT_MARGIN_TB_EMU) / EMU_PER_IN
        mb = (tf.margin_bottom if tf.margin_bottom is not None else DEFAULT_MARGIN_TB_EMU) / EMU_PER_IN
        budget_pt = max(1.0, (w_in - ml - mr) * 72)
        total_h_pt = 0.0
        for para in tf.paragraphs:
            fs = _para_font_pt(para)
            lines = _wrap_lines(para.text, fs, budget_pt)
            # leading (LINE_HEIGHT_MULT) applies BETWEEN lines; the last line has no
            # trailing leading — otherwise every single-line big number (96pt stat in a
            # 1.5" box) false-positives
            total_h_pt += (lines - 1) * fs * LINE_HEIGHT_MULT + fs
        avail_pt = max(1.0, (h_in - mt - mb) * 72)
        if total_h_pt > avail_pt * OVERFLOW_TOL:
            n += 1
            example = example or (sh.name or f"shape#{sh.shape_id}")
    return n, example


def _overlap_ratio(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ix = max(0.0, min(ax + aw, bx + bw) - max(ax, bx))
    iy = max(0.0, min(ay + ah, by + bh) - max(ay, by))
    inter = ix * iy
    if inter <= 0:
        return 0.0
    smaller = min(aw * ah, bw * bh)
    return inter / smaller if smaller > 0 else 0.0


def _is_badge(kind, box):
    """A small text box (<= BADGE_MAX_IN in both dimensions) is a badge/marker — e.g.
    arch_kit's numbered stepMarker circles, which sit ON icons/labels by design.
    Exempting badges from overlap keeps the intended-overlay pattern legal while still
    catching real caption-on-icon collisions (captions are wide, not square)."""
    return kind == "text" and box[2] <= BADGE_MAX_IN and box[3] <= BADGE_MAX_IN


def _iou(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ix = max(0.0, min(ax + aw, bx + bw) - max(ax, bx))
    iy = max(0.0, min(ay + ah, by + bh) - max(ay, by))
    inter = ix * iy
    union = aw * ah + bw * bh - inter
    return inter / union if union > 0 else 0.0


def _check_overlap(shapes, slide_area_sqin):
    boxes = []
    for sh in shapes:
        kind = _kind(sh, slide_area_sqin)
        if kind is None:
            continue
        box = _bbox_in(sh)
        if box:
            boxes.append((kind, box))
    n = 0
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            (ka, ba), (kb, bb) = boxes[i], boxes[j]
            # Two intended-overlay patterns are exempt:
            # 1. badge-sized text markers (stepMarker numbers) overlaying anything;
            # 2. a text label congruent with a picture (IoU >= 0.8) — a pill header /
            #    labeled shape, where the text is deliberately drawn ON the image.
            # Everything else counts: picture x picture, text x text, and the
            # partial picture x text collision a too-short cell height produces.
            # Full-bleed backgrounds are already excluded upstream by _kind().
            if _is_badge(ka, ba) or _is_badge(kb, bb):
                continue
            if ka != kb and _iou(ba, bb) >= 0.8:
                continue
            if _overlap_ratio(ba, bb) > OVERLAP_AREA_RATIO:
                n += 1
    return n


def _check_offcanvas(shapes, sw_in, sh_in, slide_area_sqin):
    n = 0
    for sh in shapes:
        box = _bbox_in(sh)
        if not box:
            continue
        x, y, w, h = box
        # A full-bleed background image may deliberately bleed past the edge — but only
        # exempt it if it actually COVERS the canvas (on-canvas intersection > 80% of the
        # slide). A large image parked entirely off-canvas is still a defect: area alone
        # would have waved it through.
        try:
            is_pic = sh.shape_type == MSO_SHAPE_TYPE.PICTURE
        except (ValueError, NotImplementedError):
            is_pic = False
        if is_pic:
            ix = max(0.0, min(x + w, sw_in) - max(x, 0.0))
            iy = max(0.0, min(y + h, sh_in) - max(y, 0.0))
            if (ix * iy) / slide_area_sqin > BG_AREA_RATIO:
                continue
        if (x < -CANVAS_TOL_IN or y < -CANVAS_TOL_IN or
                x + w > sw_in + CANVAS_TOL_IN or y + h > sh_in + CANVAS_TOL_IN):
            n += 1
    return n


def _footer_text(shapes):
    """Text from shapes whose box sits in the footer band only — a body paragraph that
    merely mentions 'Amazon Web Services' must not satisfy the footer check."""
    parts = []
    for sh in shapes:
        if not _has_text(sh):
            continue
        box = _bbox_in(sh)
        if box and box[1] > FOOTER_BAND_IN:
            parts.append(sh.text_frame.text)
    return " ".join(parts)


def analyze(prs):
    """Run all checks against a python-pptx Presentation object (already-opened or
    freshly built in-memory — both work, no save/reload round-trip required)."""
    sw_in = prs.slide_width / EMU_PER_IN
    sh_in = prs.slide_height / EMU_PER_IN
    slide_area = sw_in * sh_in

    findings = []   # geometry: (weight, message)
    design = []     # design: (weight, message)
    notes = []      # non-scoring caveats (e.g. GROUP coverage limitation)
    total_shapes = 0
    group_slides = []
    page_nums = []          # (slide_index, int_value) in slide order
    missing_footer = []
    missing_pagenum = []    # content slides (have footer) lacking a page number
    bad_font_count = 0
    small_font_count = 0
    placeholder_count = 0

    for si, slide in enumerate(prs.slides):
        shapes = list(_iter_shapes(slide.shapes))
        total_shapes += len(shapes)
        for sh in shapes:
            try:
                if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                    group_slides.append(si + 1)
                    break
            except (ValueError, NotImplementedError):
                pass

        n_overflow, ex = _check_overflow(shapes)
        if n_overflow:
            findings.append((min(25, 5 * n_overflow),
                              f"slide {si + 1}: {n_overflow} text box(es) overflow their bounds "
                              f"(estimated text height exceeds shape height) e.g. '{ex}'"))

        n_overlap = _check_overlap(shapes, slide_area)
        if n_overlap:
            findings.append((min(20, 5 * n_overlap),
                              f"slide {si + 1}: {n_overlap} pair(s) of overlapping picture(s)/text box(es)"))

        n_off = _check_offcanvas(shapes, sw_in, sh_in, slide_area)
        if n_off:
            findings.append((min(15, 5 * n_off),
                              f"slide {si + 1}: {n_off} shape(s) fall outside the slide canvas"))

        has_footer = COPYRIGHT_MARK in _footer_text(shapes).lower()
        if not has_footer:
            missing_footer.append(si + 1)

        slide_has_pagenum = False
        for sh in shapes:
            if not getattr(sh, "has_text_frame", False):
                continue
            t = sh.text_frame.text.strip()
            if t.isdigit() and _bbox_in(sh) and _bbox_in(sh)[1] > FOOTER_BAND_IN:
                page_nums.append((si, int(t)))
                slide_has_pagenum = True
            if PLACEHOLDER_RE.search(sh.text_frame.text):
                placeholder_count += 1
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    if r.font.name and r.font.name != FONT_NAME:
                        bad_font_count += 1
                    if r.font.size and r.font.size.pt < 8:
                        small_font_count += 1

        # a slide carrying a footer copyright is a content slide — it should also carry a
        # page number. The cover (slide 1) conventionally has the copyright but no number,
        # so it's exempt. Catches the case where SOME slides are numbered but others aren't.
        if has_footer and not slide_has_pagenum and si > 0:
            missing_pagenum.append(si + 1)

    if missing_footer:
        design.append((min(12, 3 * len(missing_footer)),
                        f"{len(missing_footer)} slide(s) missing the footer copyright line "
                        f"(slides {missing_footer}) — every slide should call addFooter()"))

    bad_seq = []
    for k in range(1, len(page_nums)):
        (_, prev), (si, cur) = page_nums[k - 1], page_nums[k]
        if cur <= prev:
            bad_seq.append((si + 1, prev, cur))
    if bad_seq:
        shown = bad_seq[0]
        design.append((min(12, 3 * len(bad_seq)),
                        f"{len(bad_seq)} page-number regression(s)/duplicate(s) — e.g. slide "
                        f"{shown[0]} shows {shown[2]} after {shown[1]} (page numbers must strictly increase)"))
    if missing_pagenum:
        design.append((min(8, 2 * len(missing_pagenum)),
                        f"{len(missing_pagenum)} content slide(s) with a footer but no page number "
                        f"(slides {missing_pagenum}) — every numbered content slide should carry one"))

    if group_slides:
        # GROUP children can't be re-projected into slide coordinates, so every geometry
        # and design check is blind to them. Rather than silently pass a deck whose
        # grouped content may be broken, this is a blocking geometry finding: it fails the
        # gate and forces a manual look (the kit itself never emits groups, so kit-built
        # decks never hit this — it guards the external-deck path content-review uses).
        findings.append((min(15, 5 * len(group_slides)),
                         f"{len(group_slides)} slide(s) contain native GROUP shapes "
                         f"(slides {group_slides}) whose children can't be inspected — "
                         f"ungroup them or review those slides by hand before shipping"))

    if bad_font_count:
        design.append((min(10, 2 * bad_font_count),
                        f"{bad_font_count} run(s) not using {FONT_NAME} — inconsistent type looks unpolished"))

    if small_font_count:
        design.append((min(8, 2 * small_font_count),
                        f"{small_font_count} run(s) below the 8pt minimum (footer floor)"))

    if placeholder_count:
        design.append((min(10, 5 * placeholder_count),
                        f"{placeholder_count} run(s) contain placeholder text (TODO/TBD/lorem/???) "
                        f"— unfinished content left in the deck"))

    geom_loss = sum(w for w, _ in findings)
    design_loss = sum(w for w, _ in design)
    score = max(0, 100 - geom_loss - design_loss)
    return {
        "score": score,
        "subscores": {"geometry": max(0, 100 - geom_loss), "design": max(0, 100 - design_loss)},
        "counts": {"slides": len(prs.slides), "shapes": total_shapes},
        "geometry_findings": len(findings),
        "findings": [f"[geometry] {m}" for _, m in findings] + [f"[design] {m}" for _, m in design],
        "notes": notes,
    }


def gate_pass(result, threshold):
    """The one gate rule (docstring: 'THE GATE'). Content-review-agent treats any
    geometry finding as Critical, so the skill-side gate must agree — score alone
    can't buy back an overflow/overlap/off-canvas defect."""
    return result["score"] >= threshold and result["geometry_findings"] == 0


def main():
    args = sys.argv[1:]
    if not args or args[0].startswith("--"):
        print(__doc__)
        return 2
    path = args[0]
    threshold = DEFAULT_THRESHOLD
    if "--threshold" in args:
        i = args.index("--threshold")
        if i + 1 >= len(args) or not args[i + 1].isdigit() or not (0 <= int(args[i + 1]) <= 100):
            print("usage: check_pptx.py <file.pptx> [--threshold N] [--json]  (N = integer 0-100)")
            return 2
        threshold = int(args[i + 1])
    as_json = "--json" in args

    try:
        prs = Presentation(path)
    except Exception as e:  # noqa: BLE001 - python-pptx raises several exception types on bad input
        print(f"❌ cannot read {path}: {e}")
        return 2

    result = analyze(prs)
    ok = gate_pass(result, threshold)

    if as_json:
        result["threshold"] = threshold
        result["pass"] = ok
        print(_json.dumps(result, ensure_ascii=False, indent=2))
        return 0 if ok else 1

    c = result["counts"]
    sub = result.get("subscores", {})
    mark = "✅" if ok else "❌"
    breakdown = f" [geometry {sub['geometry']} · design {sub['design']}]" if sub else ""
    print(f"{mark} {path}: deck score {result['score']}/100 (gate {threshold} + geometry clean){breakdown} "
          f"— slides={c['slides']} shapes={c['shapes']}")
    for m in result["findings"]:
        print(f"   • {m}")
    for m in result["notes"]:
        print(f"   ⚠ {m}")
    if not ok:
        why = ("geometry findings present (any geometry finding fails the gate regardless of score)"
               if result["geometry_findings"] else "score below threshold")
        print(f"   → QA gate failed: {why}. Fix [geometry] (text overflow, overlap, off-canvas) "
              "and [design] (footer, page numbers, Pretendard-only, no placeholder text) "
              "before delivering.")
    return 0 if ok else 1


if __name__ == "__main__":
    sys.exit(main())
