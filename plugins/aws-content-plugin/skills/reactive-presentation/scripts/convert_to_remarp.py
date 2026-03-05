#!/usr/bin/env python3
"""
PPTX/PDF → Remarp Markdown Converter.

Converts existing PPTX or PDF presentation files into Remarp markdown projects
that can be edited and built with remarp_to_slides.py.

Pipeline:  Source File → RawSlide IR → .remarp.md files → (optional) HTML build

Supported sources:
  - .pptx  — text, images, layouts, speaker notes via python-pptx
  - .pdf   — page images + text extraction via PyMuPDF

Usage:
  python3 convert_to_remarp.py input.pptx -o output_dir/
  python3 convert_to_remarp.py input.pdf  -o output_dir/ --build
  python3 convert_to_remarp.py input.pptx -o output_dir/ --block-size 15 --lang ko
"""

import argparse
import os
import re
import shutil
import subprocess
import sys
import time
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from datetime import date
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple


# ---------------------------------------------------------------------------
# Data Model (Intermediate Representation)
# ---------------------------------------------------------------------------

@dataclass
class ImageRef:
    """Reference to an extracted image file."""
    path: str           # Relative path from project root (e.g., 'assets/slide-01-img.png')
    alt_text: str = ''
    width_pct: int = 80  # Display width as percentage


@dataclass
class TextBox:
    """A positioned text box from the source slide."""
    text: str
    left_pct: float = 0.0
    top_pct: float = 0.0
    width_pct: float = 100.0
    height_pct: float = 100.0
    font_size: float = 0.0
    is_title: bool = False


@dataclass
class RawSlide:
    """Intermediate representation for a single slide extracted from source."""
    index: int
    title: str = ''
    body_text: str = ''
    speaker_notes: str = ''
    layout_name: str = ''
    background_image: Optional[str] = None
    images: List[ImageRef] = field(default_factory=list)
    has_table: bool = False
    has_chart: bool = False
    bullet_items: List[str] = field(default_factory=list)
    text_boxes: List[TextBox] = field(default_factory=list)


@dataclass
class Block:
    """A group of slides forming one presentation block."""
    name: str
    title: str
    slides: List[RawSlide] = field(default_factory=list)
    duration: int = 25  # minutes


# ---------------------------------------------------------------------------
# Layout → Remarp Type Mapping
# ---------------------------------------------------------------------------

# PPTX layout names → Remarp @type directives
LAYOUT_TYPE_MAP = {
    'Title Slide': 'cover',
    'Section Header': 'title',
    'Title and Content': 'content',
    'Two Content': 'content',       # + @layout: two-column
    'Comparison': 'compare',
    'Title Only': 'content',
    'Blank': 'content',
    'Content with Caption': 'content',
    'Picture with Caption': 'content',
}

# Layouts that imply two-column layout
TWO_COLUMN_LAYOUTS = {'Two Content', 'Comparison'}


# ---------------------------------------------------------------------------
# Base Converter
# ---------------------------------------------------------------------------

class SourceConverter(ABC):
    """Base class for source file conversion."""

    def __init__(self, input_path: str, output_dir: str):
        self.input_path = Path(input_path)
        self.output_dir = Path(output_dir)
        self.assets_dir = self.output_dir / 'assets'
        self.assets_dir.mkdir(parents=True, exist_ok=True)

    @abstractmethod
    def extract_slides(self) -> List[RawSlide]:
        """Extract slides from source file into RawSlide IR."""
        ...


# ---------------------------------------------------------------------------
# PPTX Converter
# ---------------------------------------------------------------------------

class PPTXConverter(SourceConverter):
    """Convert PPTX files to RawSlide list using python-pptx."""

    def __init__(self, input_path: str, output_dir: str):
        super().__init__(input_path, output_dir)
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            print("Error: python-pptx is required for PPTX conversion.")
            print("Install: pip install python-pptx")
            sys.exit(1)
        self._MSO_SHAPE_TYPE = MSO_SHAPE_TYPE
        self.prs = Presentation(str(self.input_path))
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height

    def extract_slides(self) -> List[RawSlide]:
        slides = []
        for i, slide in enumerate(self.prs.slides):
            raw = RawSlide(index=i)
            raw.layout_name = slide.slide_layout.name if slide.slide_layout else ''

            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                raw.speaker_notes = slide.notes_slide.notes_text_frame.text.strip()

            for shape in slide.shapes:
                if shape.has_text_frame:
                    self._extract_text(shape, raw)
                if shape.shape_type == self._MSO_SHAPE_TYPE.PICTURE:
                    self._extract_image(shape, raw, i)
                if shape.has_table:
                    raw.has_table = True
                    self._extract_table(shape, raw)
                if shape.has_chart:
                    raw.has_chart = True

            slides.append(raw)
        return slides

    def _extract_text(self, shape, raw: RawSlide):
        """Extract text from a shape's text frame, converting to markdown."""
        tf = shape.text_frame
        is_title = shape.is_placeholder and hasattr(shape, 'placeholder_format')
        if is_title:
            ph_type = shape.placeholder_format.type
            # Title placeholder types: TITLE=15, CENTER_TITLE=3, SUBTITLE=4
            if ph_type is not None and ph_type in (15, 3):
                raw.title = self._text_frame_to_plain(tf)
                return
            if ph_type is not None and ph_type == 4:
                # Subtitle — append to body
                subtitle = self._text_frame_to_plain(tf)
                if subtitle:
                    raw.body_text = subtitle + ('\n\n' + raw.body_text if raw.body_text else '')
                return

        md = self._text_frame_to_markdown(tf)
        if md.strip():
            if raw.body_text:
                raw.body_text += '\n\n' + md
            else:
                raw.body_text = md

    def _text_frame_to_plain(self, text_frame) -> str:
        """Extract plain text from a TextFrame."""
        parts = []
        for para in text_frame.paragraphs:
            line = ''.join(run.text for run in para.runs)
            if line.strip():
                parts.append(line.strip())
        return ' '.join(parts)

    def _text_frame_to_markdown(self, text_frame) -> str:
        """Convert TextFrame to markdown with bullet levels and formatting."""
        lines = []
        for para in text_frame.paragraphs:
            level = para.level if para.level else 0

            parts = []
            for run in para.runs:
                text = run.text
                if not text:
                    continue
                if run.font.bold:
                    text = f'**{text}**'
                if run.font.italic:
                    text = f'*{text}*'
                parts.append(text)

            line_text = ''.join(parts).strip()
            if not line_text:
                continue

            # Build bullet prefix based on indentation level
            if level > 0:
                prefix = '  ' * (level - 1) + '  - '
            else:
                prefix = '- '

            lines.append(f'{prefix}{line_text}')

        return '\n'.join(lines)

    def _extract_image(self, shape, raw: RawSlide, slide_idx: int):
        """Save slide image to assets/ and record reference."""
        try:
            image = shape.image
            content_type = image.content_type or 'image/png'
            ext = content_type.split('/')[-1].replace('jpeg', 'jpg')
            # Sanitize shape name for filename
            safe_name = re.sub(r'[^\w\-]', '_', shape.name or f'img{shape.shape_id}')
            filename = f'slide-{slide_idx:02d}-{safe_name}.{ext}'

            image_path = self.assets_dir / filename
            with open(image_path, 'wb') as f:
                f.write(image.blob)

            # Determine if this is a background image (covers >80% of slide)
            is_bg = (
                shape.width and shape.height
                and shape.width > self.slide_width * 0.8
                and shape.height > self.slide_height * 0.8
            )

            rel_path = f'assets/{filename}'
            if is_bg:
                raw.background_image = rel_path
            else:
                raw.images.append(ImageRef(
                    path=rel_path,
                    alt_text=shape.name or '',
                    width_pct=min(80, int((shape.width / self.slide_width) * 100)) if shape.width else 80
                ))
        except Exception:
            pass  # Skip unreadable images

    def _extract_table(self, shape, raw: RawSlide):
        """Convert table to markdown."""
        table = shape.table
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)

        if not rows:
            return

        # Build markdown table
        md_lines = []
        # Header row
        md_lines.append('| ' + ' | '.join(rows[0]) + ' |')
        md_lines.append('| ' + ' | '.join(['---'] * len(rows[0])) + ' |')
        # Data rows
        for row in rows[1:]:
            # Pad row to header length
            while len(row) < len(rows[0]):
                row.append('')
            md_lines.append('| ' + ' | '.join(row) + ' |')

        table_md = '\n'.join(md_lines)
        if raw.body_text:
            raw.body_text += '\n\n' + table_md
        else:
            raw.body_text = table_md

    def extract_theme(self):
        """Extract PPTX theme using the existing extract_pptx_theme.py script."""
        theme_dir = self.output_dir / '_theme' / self.input_path.stem
        theme_dir.mkdir(parents=True, exist_ok=True)

        manifest_path = theme_dir / 'theme-manifest.json'
        if manifest_path.exists():
            src_mtime = self.input_path.stat().st_mtime
            manifest_mtime = manifest_path.stat().st_mtime
            if manifest_mtime > src_mtime:
                print(f"  Theme cache valid: {theme_dir.name}/")
                return theme_dir

        script_dir = Path(__file__).parent
        extract_script = script_dir / 'extract_pptx_theme.py'

        if not extract_script.exists():
            print(f"  Warning: extract_pptx_theme.py not found at {extract_script}")
            return None

        print(f"  Extracting theme from {self.input_path.name}...")
        result = subprocess.run(
            [sys.executable, str(extract_script), str(self.input_path), '-o', str(theme_dir)],
            capture_output=True, text=True
        )
        if result.returncode != 0:
            print(f"  Warning: Theme extraction failed: {result.stderr.strip()}")
            return None

        print(f"  Theme extracted to {theme_dir.name}/")
        return theme_dir


# ---------------------------------------------------------------------------
# PDF Converter
# ---------------------------------------------------------------------------

class PDFConverter(SourceConverter):
    """Convert PDF files to RawSlide list using PyMuPDF (fitz)."""

    def __init__(self, input_path: str, output_dir: str):
        super().__init__(input_path, output_dir)
        try:
            import fitz  # noqa: F401
        except ImportError:
            print("Error: PyMuPDF is required for PDF conversion.")
            print("Install: pip install PyMuPDF")
            sys.exit(1)

    def extract_slides(self) -> List[RawSlide]:
        import fitz

        doc = fitz.open(str(self.input_path))
        slides = []

        for i, page in enumerate(doc):
            raw = RawSlide(index=i)

            # Render page as high-res image for background
            pix = page.get_pixmap(dpi=200)
            img_filename = f'page-{i:02d}.png'
            img_path = self.assets_dir / img_filename
            pix.save(str(img_path))
            raw.background_image = f'assets/{img_filename}'

            # Extract structured text
            try:
                blocks = page.get_text("dict")["blocks"]
                raw.title, raw.body_text = self._classify_text_blocks(blocks)
            except Exception:
                raw.title = f'Slide {i + 1}'
                raw.body_text = page.get_text("text").strip()

            # Layout hint for PDF: all slides are content type
            raw.layout_name = 'Title Slide' if i == 0 else 'Content'

            slides.append(raw)

        doc.close()
        return slides

    def _classify_text_blocks(self, blocks: List[Dict[str, Any]]) -> Tuple[str, str]:
        """Classify text blocks into title and body based on font size."""
        text_blocks = [b for b in blocks if b.get('type') == 0]  # text blocks only
        if not text_blocks:
            return '', ''

        # Find maximum font size across all spans
        all_spans = []
        for block in text_blocks:
            for line in block.get('lines', []):
                for span in line.get('spans', []):
                    if span.get('text', '').strip():
                        all_spans.append(span)

        if not all_spans:
            return '', ''

        max_size = max(s['size'] for s in all_spans)

        title_parts = []
        body_parts = []

        for block in text_blocks:
            for line in block.get('lines', []):
                for span in line.get('spans', []):
                    text = span.get('text', '').strip()
                    if not text:
                        continue
                    if span['size'] >= max_size * 0.85:
                        title_parts.append(text)
                    else:
                        body_parts.append(text)

        title = ' '.join(title_parts).strip()
        body = '\n'.join(body_parts).strip()
        return title, body


# ---------------------------------------------------------------------------
# Remarp Generator
# ---------------------------------------------------------------------------

class RemarpGenerator:
    """Generate Remarp markdown project from RawSlide list."""

    def __init__(self, output_dir: str, lang: str = 'ko',
                 block_size: int = 15, theme_dir: Optional[Path] = None,
                 source_stem: str = 'presentation'):
        self.output_dir = Path(output_dir)
        self.lang = lang
        self.block_size = block_size
        self.theme_dir = theme_dir
        self.source_stem = source_stem

    def generate(self, slides: List[RawSlide]) -> List[str]:
        """Generate Remarp project files. Returns list of created file paths."""
        if not slides:
            print("  Warning: No slides to convert.")
            return []

        blocks = self._split_into_blocks(slides)
        created = []

        # Generate _presentation.md
        pres_path = self._generate_presentation_md(blocks, slides[0].title)
        created.append(str(pres_path))

        # Generate block files
        for block in blocks:
            block_path = self._generate_block_md(block)
            created.append(str(block_path))

        return created

    def _split_into_blocks(self, slides: List[RawSlide]) -> List[Block]:
        """Split slides into blocks using Section Headers or uniform sizing."""
        # Find section header positions
        section_indices = [
            i for i, s in enumerate(slides)
            if s.layout_name in ('Section Header', 'Section')
        ]

        if len(section_indices) >= 2:
            return self._split_by_sections(slides, section_indices)
        else:
            return self._split_uniform(slides)

    def _split_by_sections(self, slides: List[RawSlide],
                           section_indices: List[int]) -> List[Block]:
        """Split slides at Section Header boundaries."""
        blocks = []
        # Add boundaries: start and end
        boundaries = [0] + section_indices + [len(slides)]
        # Remove duplicates and sort
        boundaries = sorted(set(boundaries))

        for idx, (start, end) in enumerate(zip(boundaries, boundaries[1:])):
            block_slides = slides[start:end]
            if not block_slides:
                continue

            # Derive block title from section header or first slide
            title = block_slides[0].title or f'Part {idx + 1}'
            name = f'{idx + 1:02d}-{self._slugify(title)}'

            blocks.append(Block(
                name=name,
                title=title,
                slides=block_slides,
                duration=max(15, len(block_slides) * 2)  # ~2 min per slide
            ))

        return blocks

    def _split_uniform(self, slides: List[RawSlide]) -> List[Block]:
        """Split slides into uniform blocks of block_size."""
        blocks = []
        for idx, start in enumerate(range(0, len(slides), self.block_size)):
            block_slides = slides[start:start + self.block_size]
            title = block_slides[0].title or f'Part {idx + 1}'
            name = f'{idx + 1:02d}-{self._slugify(title)}'

            blocks.append(Block(
                name=name,
                title=title,
                slides=block_slides,
                duration=max(15, len(block_slides) * 2)
            ))

        return blocks

    def _generate_presentation_md(self, blocks: List[Block],
                                  first_title: str) -> Path:
        """Generate _presentation.md with global config."""
        title = first_title or self.source_stem
        today = date.today().isoformat()

        lines = ['---', 'remarp: true', 'version: 1']
        lines.append(f'title: "{title}"')
        lines.append(f'date: {today}')
        lines.append(f'lang: {self.lang}')
        lines.append('')

        # Theme section
        if self.theme_dir and self.theme_dir.exists():
            rel_theme = os.path.relpath(self.theme_dir, self.output_dir)
            lines.append('theme:')
            lines.append(f'  source: "./{rel_theme}"')
            lines.append('  footer: auto')
            lines.append('  logo: auto')
        else:
            lines.append('theme:')
            lines.append('  footer: ""')

        lines.append('')

        # Blocks section
        lines.append('blocks:')
        for block in blocks:
            lines.append(f'  - name: {block.name}')
            lines.append(f'    title: "{block.title}"')
            lines.append(f'    duration: {block.duration}')

        lines.append('---')

        out_path = self.output_dir / '_presentation.md'
        out_path.write_text('\n'.join(lines), encoding='utf-8')
        return out_path

    def _generate_block_md(self, block: Block) -> Path:
        """Generate a single block .md file."""
        lines = ['---', 'remarp: true', f'block: {block.name}', '---', '']

        for i, slide in enumerate(block.slides):
            is_first_slide = (i == 0 and block.slides[0].index == 0)
            slide_md = self._slide_to_remarp(slide, is_first_slide)
            lines.append(slide_md)

        out_path = self.output_dir / f'{block.name}.md'
        out_path.write_text('\n'.join(lines), encoding='utf-8')
        return out_path

    def _slide_to_remarp(self, slide: RawSlide, is_first: bool) -> str:
        """Convert a single RawSlide to Remarp slide markdown."""
        parts = ['---']

        # Directives
        layout = slide.layout_name
        remarp_type = LAYOUT_TYPE_MAP.get(layout, 'content')

        if is_first and remarp_type in ('cover', 'content'):
            remarp_type = 'cover'
        elif layout == 'Section Header':
            remarp_type = 'title'

        parts.append(f'@type: {remarp_type}')

        if layout in TWO_COLUMN_LAYOUTS:
            parts.append('@layout: two-column')

        if slide.background_image:
            parts.append(f'@background: {slide.background_image}')

        parts.append('')  # Blank line separating directives from content

        # Title
        if slide.title:
            prefix = '# ' if is_first or remarp_type == 'cover' else '## '
            parts.append(f'{prefix}{slide.title}')

        # Body text
        if slide.body_text:
            parts.append('')
            parts.append(slide.body_text)

        # Inline images (non-background)
        for img in slide.images:
            parts.append(f'\n@img: {img.path} center {img.width_pct}%')

        # Chart/table indicators as comments
        if slide.has_chart:
            parts.append('\n<!-- TODO: Chart detected — recreate with Canvas DSL or replace with data -->')

        # Speaker notes
        if slide.speaker_notes:
            parts.append('')
            parts.append(':::notes')
            parts.append(slide.speaker_notes)
            parts.append(':::')

        parts.append('')  # Trailing newline
        return '\n'.join(parts)

    @staticmethod
    def _slugify(text: str) -> str:
        """Convert text to URL-safe slug."""
        # Remove non-alphanumeric (keep Korean, spaces, hyphens)
        slug = re.sub(r'[^\w\s\-]', '', text, flags=re.UNICODE)
        # Replace whitespace with hyphens
        slug = re.sub(r'\s+', '-', slug.strip())
        # Lowercase and truncate
        slug = slug.lower()[:40].rstrip('-')
        return slug or 'untitled'


# ---------------------------------------------------------------------------
# Orchestrator
# ---------------------------------------------------------------------------

def convert_to_remarp(input_path: str, output_dir: str, lang: str = 'ko',
                      block_size: int = 15, build: bool = False,
                      force: bool = False) -> Optional[List[str]]:
    """Main conversion pipeline.

    Returns list of created file paths, or None on failure.
    """
    src = Path(input_path)
    out = Path(output_dir)
    ext = src.suffix.lower()

    if not src.exists():
        print(f"Error: Source file not found: {src}")
        return None

    if ext not in ('.pptx', '.pdf'):
        print(f"Error: Unsupported format '{ext}'. Use .pptx or .pdf")
        if ext == '.ppt':
            print("Hint: .ppt (legacy) is not supported. Save as .pptx in PowerPoint first.")
        return None

    # Overwrite protection
    if out.exists() and any(out.glob('*.md')):
        if not force:
            print(f"Error: {out} already contains .md files.")
            print("Use --force to overwrite (creates timestamped .bak backup)")
            return None
        # Create backup
        backup_dir = out.parent / f'{out.name}.bak.{int(time.time())}'
        print(f"  Backing up existing files to {backup_dir.name}/")
        shutil.copytree(out, backup_dir)
        # Remove existing .md files only (keep assets if they exist)
        for md_file in out.glob('*.md'):
            md_file.unlink()

    out.mkdir(parents=True, exist_ok=True)

    print(f"Converting: {src.name} → {out}/")
    print(f"  Format: {ext.upper().lstrip('.')}")
    print(f"  Language: {lang}")

    # Step 1: Create converter
    if ext == '.pptx':
        converter = PPTXConverter(input_path, output_dir)
    else:
        converter = PDFConverter(input_path, output_dir)

    # Step 2: Extract theme (PPTX only)
    theme_dir = None
    if ext == '.pptx':
        theme_dir = converter.extract_theme()

    # Step 3: Extract slides
    print("  Extracting slides...")
    raw_slides = converter.extract_slides()
    print(f"  Found {len(raw_slides)} slides")

    if not raw_slides:
        print("Error: No slides found in source file.")
        return None

    # Step 4: Generate Remarp project
    print("  Generating Remarp markdown...")
    generator = RemarpGenerator(
        output_dir=output_dir,
        lang=lang,
        block_size=block_size,
        theme_dir=theme_dir,
        source_stem=src.stem
    )
    created_files = generator.generate(raw_slides)

    # Summary
    block_files = [f for f in created_files if not f.endswith('_presentation.md')]
    asset_count = len(list((out / 'assets').glob('*'))) if (out / 'assets').exists() else 0
    print(f"\n  Created {len(created_files)} files:")
    print(f"    _presentation.md (global config)")
    for bf in block_files:
        print(f"    {Path(bf).name}")
    print(f"    assets/ ({asset_count} files)")
    if theme_dir:
        print(f"    _theme/{theme_dir.name}/ (extracted theme)")

    # Step 5: Optional HTML build
    if build:
        print("\n  Building HTML...")
        script_dir = Path(__file__).parent
        build_script = script_dir / 'remarp_to_slides.py'
        if build_script.exists():
            result = subprocess.run(
                [sys.executable, str(build_script), 'build', str(out)],
                capture_output=True, text=True
            )
            if result.returncode == 0:
                print("  HTML build complete.")
                if result.stdout:
                    # Print last few lines of build output
                    for line in result.stdout.strip().split('\n')[-5:]:
                        print(f"    {line}")
            else:
                print(f"  Warning: HTML build failed: {result.stderr.strip()}")
        else:
            print(f"  Warning: remarp_to_slides.py not found at {build_script}")

    print(f"\nDone! Edit the .md files in {out}/ then build with:")
    print(f"  python3 {Path(__file__).parent}/remarp_to_slides.py build {out}/")

    return created_files


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description='Convert PPTX/PDF to Remarp markdown project.',
        epilog='Examples:\n'
               '  %(prog)s slides.pptx -o ./converted/\n'
               '  %(prog)s slides.pdf -o ./converted/ --build\n'
               '  %(prog)s template.pptx -o ./project/ --block-size 10 --lang en\n',
        formatter_class=argparse.RawDescriptionHelpFormatter
    )
    parser.add_argument('input', help='Source PPTX or PDF file')
    parser.add_argument('-o', '--output', required=True,
                        help='Output directory for Remarp project')
    parser.add_argument('--build', action='store_true',
                        help='Also build HTML after conversion')
    parser.add_argument('--block-size', type=int, default=15,
                        help='Max slides per block (default: 15, ignored if Section Headers found)')
    parser.add_argument('--lang', choices=['ko', 'en'], default='ko',
                        help='Presentation language (default: ko)')
    parser.add_argument('--force', action='store_true',
                        help='Overwrite existing output (creates .bak backup)')

    args = parser.parse_args()

    result = convert_to_remarp(
        input_path=args.input,
        output_dir=args.output,
        lang=args.lang,
        block_size=args.block_size,
        build=args.build,
        force=args.force,
    )

    sys.exit(0 if result else 1)


if __name__ == '__main__':
    main()
