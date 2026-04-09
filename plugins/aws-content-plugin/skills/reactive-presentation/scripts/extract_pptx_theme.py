#!/usr/bin/env python3
"""
PPTX Theme Extraction Script

Extracts theme colors, fonts, backgrounds, logos, footer information,
and slide master template metadata from PowerPoint presentations.
Generates CSS overrides for dark-theme presentations.

The ``slide_master`` section in the manifest captures the common frame
(header, footer, decorative elements, color palette, slide size) that
applies uniformly across all slides — i.e. the PPT slide master, not
individual slide content layouts.

Requires: python-pptx >= 1.0.0, lxml
"""

import argparse
import json
import os
import sys
from pathlib import Path
from typing import Any, Optional

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

# XML namespaces for Office Open XML
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

# Standard 16:9 slide dimensions in EMU
SLIDE_WIDTH_EMU = 12192000
SLIDE_HEIGHT_EMU = 6858000

# Placeholder type constants
PH_TYPE_FOOTER = 15
PH_TYPE_SLIDE_NUMBER = 13
PH_TYPE_DATE = 16


def emu_to_percent(emu_value: int, dimension: int) -> float:
    """Convert EMU to percentage of slide dimension."""
    return round((emu_value / dimension) * 100, 4)


def hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def rgb_to_luminance(r: int, g: int, b: int) -> float:
    """Calculate relative luminance of a color."""
    def channel_luminance(c: int) -> float:
        c = c / 255
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    return 0.2126 * channel_luminance(r) + 0.7152 * channel_luminance(g) + 0.0722 * channel_luminance(b)


class ThemeExtractor:
    """Extracts theme information from a PPTX file."""

    def __init__(self, pptx_path: str, master_index: int = 0):
        """
        Initialize extractor with PPTX file.

        Args:
            pptx_path: Path to the PPTX file
            master_index: Index of slide master to use (default: 0)
        """
        self.pptx_path = Path(pptx_path)
        if not self.pptx_path.exists():
            raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

        self.prs = Presentation(pptx_path)
        self.master_index = master_index

        if master_index >= len(self.prs.slide_masters):
            raise IndexError(f"Master index {master_index} out of range. "
                           f"Only {len(self.prs.slide_masters)} masters available.")

        self.master = self.prs.slide_masters[master_index]
        self.theme_xml = self._get_theme_xml()
        self.theme_name = self._get_theme_name()

    def _get_theme_xml(self) -> Optional[etree._Element]:
        """Get the theme XML from the slide master's relationships."""
        for rel in self.master.part.rels.values():
            if 'theme' in rel.reltype:
                theme_part = rel.target_part
                return etree.fromstring(theme_part.blob)
        return None

    def _get_theme_name(self) -> str:
        """Extract theme name from theme XML."""
        if self.theme_xml is not None:
            name = self.theme_xml.get('name')
            if name:
                return name
        return "Unknown"

    def extract_colors(self) -> dict[str, str]:
        """
        Extract color scheme from theme XML.

        Returns:
            Dictionary mapping color names to hex values.
        """
        if hasattr(self, '_theme_colors') and self._theme_colors:
            return dict(self._theme_colors)

        colors = {}
        if self.theme_xml is None:
            return colors

        color_scheme = self.theme_xml.find('.//a:clrScheme', NSMAP)
        if color_scheme is None:
            return colors

        color_names = ['dk1', 'lt1', 'dk2', 'lt2',
                      'accent1', 'accent2', 'accent3', 'accent4',
                      'accent5', 'accent6', 'hlink', 'folHlink']

        for name in color_names:
            elem = color_scheme.find(f'a:{name}', NSMAP)
            if elem is not None:
                # Check for srgbClr (explicit RGB)
                srgb = elem.find('a:srgbClr', NSMAP)
                if srgb is not None:
                    colors[name] = f"#{srgb.get('val')}"
                    continue

                # Check for sysClr (system color with lastClr fallback)
                sys_clr = elem.find('a:sysClr', NSMAP)
                if sys_clr is not None:
                    last_clr = sys_clr.get('lastClr')
                    if last_clr:
                        colors[name] = f"#{last_clr}"
                    else:
                        # Fallback based on system color name
                        sys_val = sys_clr.get('val', '')
                        if 'windowText' in sys_val:
                            colors[name] = '#000000'
                        elif 'window' in sys_val:
                            colors[name] = '#FFFFFF'

        self._theme_colors = colors
        return colors

    def extract_fonts(self) -> dict[str, str]:
        """
        Extract font scheme from theme XML.

        Returns:
            Dictionary with 'heading' and 'body' font names.
        """
        fonts = {'heading': 'Calibri Light', 'body': 'Calibri'}
        if self.theme_xml is None:
            return fonts

        font_scheme = self.theme_xml.find('.//a:fontScheme', NSMAP)
        if font_scheme is None:
            return fonts

        # Major font (headings)
        major = font_scheme.find('a:majorFont/a:latin', NSMAP)
        if major is not None:
            typeface = major.get('typeface')
            if typeface:
                fonts['heading'] = typeface

        # Minor font (body)
        minor = font_scheme.find('a:minorFont/a:latin', NSMAP)
        if minor is not None:
            typeface = minor.get('typeface')
            if typeface:
                fonts['body'] = typeface

        return fonts

    def extract_backgrounds(self, output_dir: Optional[Path] = None) -> dict[str, Any]:
        """
        Extract background information from master and key layouts.

        Args:
            output_dir: Directory to save extracted background images.

        Returns:
            Dictionary with background type and properties.
        """
        backgrounds = {
            'master': self._extract_background_from_element(self.master, output_dir),
            'layouts': {}
        }

        # Check key layouts
        key_layout_names = ['Title Slide', 'Title and Content', 'Section Header',
                          'Two Content', 'Blank']

        for layout in self.master.slide_layouts:
            if layout.name in key_layout_names:
                bg_info = self._extract_background_from_element(layout, output_dir)
                if bg_info['type'] != 'inherited':
                    backgrounds['layouts'][layout.name] = bg_info

        return backgrounds

    def _extract_background_from_element(self, element, output_dir: Optional[Path] = None) -> dict[str, Any]:
        """Extract background info from a slide master or layout.

        When output_dir is provided and fill is a picture, the image blob
        is saved to output_dir/images/background.{ext}.
        """
        bg_info = {'type': 'inherited'}

        try:
            background = element.background
            if background is None:
                return bg_info

            fill = background.fill
            if fill is None:
                return bg_info

            fill_type = str(fill.type) if fill.type else 'none'

            if 'SOLID' in fill_type:
                bg_info['type'] = 'solid'
                try:
                    fore_color = fill.fore_color
                    if fore_color and fore_color.rgb:
                        bg_info['color'] = f"#{fore_color.rgb}"
                except Exception:
                    pass

            elif 'PICTURE' in fill_type:
                bg_info['type'] = 'picture'
                # Extract actual image blob from the background fill
                try:
                    bg_elem = background._element
                    blip = bg_elem.find('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                    if blip is not None:
                        r_embed = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        if r_embed and hasattr(element, 'part'):
                            image_part = element.part.related_parts.get(r_embed)
                            if image_part is not None:
                                blob = image_part.blob
                                content_type = image_part.content_type
                                ext = content_type.split('/')[-1]
                                if ext == 'jpeg':
                                    ext = 'jpg'
                                bg_info['image_filename'] = f'background.{ext}'

                                # Save blob if output directory provided
                                if output_dir:
                                    images_dir = output_dir / 'images'
                                    images_dir.mkdir(parents=True, exist_ok=True)
                                    img_path = images_dir / f'background.{ext}'
                                    with open(img_path, 'wb') as f:
                                        f.write(blob)
                except Exception as e:
                    print(f"Warning: Could not extract background image: {e}",
                          file=sys.stderr)

            elif 'GRADIENT' in fill_type:
                bg_info['type'] = 'gradient'
                try:
                    bg_info['stops'] = []
                    for stop in fill.gradient_stops:
                        stop_info = {
                            'position': stop.position,
                        }
                        if stop.color and stop.color.rgb:
                            stop_info['color'] = f"#{stop.color.rgb}"
                        bg_info['stops'].append(stop_info)
                except Exception:
                    pass

            elif 'BACKGROUND' in fill_type or fill_type == 'none':
                bg_info['type'] = 'inherited'

        except Exception:
            pass

        return bg_info

    def extract_logos(self, output_dir: Path) -> list[dict[str, Any]]:
        """
        Extract logo images from master shapes.

        Filters by size heuristic: images smaller than 20% of slide width
        are likely logos.

        Args:
            output_dir: Directory to save extracted images

        Returns:
            List of logo information dictionaries.
        """
        logos = []
        images_dir = output_dir / 'images'
        images_dir.mkdir(parents=True, exist_ok=True)

        logo_threshold = SLIDE_WIDTH_EMU * 0.20  # 20% of slide width

        for shape in self.master.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                width = shape.width
                height = shape.height

                # Filter: small images are likely logos
                if width < logo_threshold and height < logo_threshold:
                    try:
                        image = shape.image
                        blob = image.blob
                        content_type = image.content_type
                        ext = content_type.split('/')[-1]
                        if ext == 'jpeg':
                            ext = 'jpg'

                        # Save image
                        filename = f"logo_{len(logos) + 1}.{ext}"
                        image_path = images_dir / filename
                        with open(image_path, 'wb') as f:
                            f.write(blob)

                        logo_info = {
                            'name': shape.name,
                            'filename': filename,
                            'content_type': content_type,
                            'position': {
                                'left_emu': shape.left,
                                'top_emu': shape.top,
                                'left_percent': emu_to_percent(shape.left, SLIDE_WIDTH_EMU),
                                'top_percent': emu_to_percent(shape.top, SLIDE_HEIGHT_EMU),
                            },
                            'size': {
                                'width_emu': width,
                                'height_emu': height,
                                'width_percent': emu_to_percent(width, SLIDE_WIDTH_EMU),
                                'height_percent': emu_to_percent(height, SLIDE_HEIGHT_EMU),
                            }
                        }

                        # Look for nearby text boxes (company name)
                        logo_info['nearby_text'] = self._find_nearby_text(shape)

                        logos.append(logo_info)

                    except Exception as e:
                        print(f"Warning: Could not extract image from {shape.name}: {e}",
                              file=sys.stderr)

        return logos

    def _find_nearby_text(self, logo_shape) -> Optional[str]:
        """Find text boxes near the logo position."""
        logo_left = logo_shape.left
        logo_top = logo_shape.top
        logo_right = logo_left + logo_shape.width
        proximity_threshold = Emu(1000000)  # ~1 inch in EMU

        for shape in self.master.shapes:
            if hasattr(shape, 'text_frame') and shape != logo_shape:
                try:
                    shape_left = shape.left
                    shape_top = shape.top

                    # Check if shape is horizontally near the logo
                    horizontal_near = abs(shape_left - logo_right) < proximity_threshold
                    vertical_aligned = abs(shape_top - logo_top) < proximity_threshold

                    if horizontal_near and vertical_aligned:
                        text = shape.text_frame.text.strip()
                        if text:
                            return text
                except Exception:
                    pass

        return None

    def extract_footer(self) -> Optional[dict[str, Any]]:
        """Extract footer placeholder information."""
        return self._extract_placeholder(PH_TYPE_FOOTER, 3)

    def extract_slide_number(self) -> Optional[dict[str, Any]]:
        """Extract slide number placeholder information."""
        return self._extract_placeholder(PH_TYPE_SLIDE_NUMBER, 4)

    def extract_date(self) -> Optional[dict[str, Any]]:
        """Extract date placeholder information."""
        return self._extract_placeholder(PH_TYPE_DATE, 2)

    def _extract_placeholder(self, ph_type: int, ph_idx: int) -> Optional[dict[str, Any]]:
        """Extract placeholder by type and index."""
        # Check master shapes
        for shape in self.master.shapes:
            if shape.is_placeholder:
                ph = shape.placeholder_format
                if ph.idx == ph_idx or (ph.type and ph.type.value == ph_type):
                    return self._get_placeholder_info(shape, ph)

        # Check layouts
        for layout in self.master.slide_layouts:
            for shape in layout.shapes:
                if shape.is_placeholder:
                    ph = shape.placeholder_format
                    if ph.idx == ph_idx or (ph.type and ph.type.value == ph_type):
                        return self._get_placeholder_info(shape, ph)

        return None

    def _get_placeholder_info(self, shape, ph) -> dict[str, Any]:
        """Get placeholder information."""
        info = {
            'name': shape.name,
            'idx': ph.idx,
            'type': ph.type.value if ph.type else None,
            'type_name': str(ph.type) if ph.type else None,
            'position': {
                'left_emu': shape.left,
                'top_emu': shape.top,
                'left_percent': emu_to_percent(shape.left, SLIDE_WIDTH_EMU),
                'top_percent': emu_to_percent(shape.top, SLIDE_HEIGHT_EMU),
            },
            'size': {
                'width_emu': shape.width,
                'height_emu': shape.height,
                'width_percent': emu_to_percent(shape.width, SLIDE_WIDTH_EMU),
                'height_percent': emu_to_percent(shape.height, SLIDE_HEIGHT_EMU),
            }
        }

        # Extract text if available
        try:
            if hasattr(shape, 'text_frame'):
                info['text'] = shape.text_frame.text
        except Exception:
            pass

        return info

    def extract_master_texts(self) -> list[dict[str, Any]]:
        """Extract non-placeholder text shapes from slide master.

        These often contain copyright notices, event names, confidentiality
        statements, and other branding text placed directly on the master.
        Deduplicates identical text across shapes.
        """
        texts = []
        seen = set()

        for shape in self.master.shapes:
            if shape.is_placeholder:
                continue
            if not hasattr(shape, 'text_frame'):
                continue
            try:
                text = shape.text_frame.text.strip()
                if not text or text in seen:
                    continue
                seen.add(text)

                info = {
                    'text': text,
                    'shape_name': shape.name,
                    'position': {
                        'left_percent': emu_to_percent(shape.left, SLIDE_WIDTH_EMU),
                        'top_percent': emu_to_percent(shape.top, SLIDE_HEIGHT_EMU),
                    },
                    'size': {
                        'width_percent': emu_to_percent(shape.width, SLIDE_WIDTH_EMU),
                        'height_percent': emu_to_percent(shape.height, SLIDE_HEIGHT_EMU),
                    },
                    'is_footer_area': emu_to_percent(shape.top, SLIDE_HEIGHT_EMU) > 85,
                }
                texts.append(info)
            except Exception:
                pass

        return texts

    # ------------------------------------------------------------------
    # Slide-master template metadata (common frame across all slides)
    # ------------------------------------------------------------------

    def extract_slide_size(self) -> dict[str, Any]:
        """Extract slide dimensions and compute aspect ratio.

        Returns e.g. ``{"width_emu": 12192000, "height_emu": 6858000,
        "width_px": 1280, "height_px": 720, "aspect_ratio": "16:9"}``.
        """
        w = self.prs.slide_width or SLIDE_WIDTH_EMU
        h = self.prs.slide_height or SLIDE_HEIGHT_EMU

        from math import gcd
        g = gcd(int(w), int(h))
        ratio_w, ratio_h = int(w) // g, int(h) // g

        # Normalise common ratios that gcd misses due to rounding
        KNOWN_RATIOS = {
            (16, 9): "16:9", (4, 3): "4:3", (16, 10): "16:10",
            (3, 2): "3:2", (5, 4): "5:4",
        }
        ratio_str = KNOWN_RATIOS.get((ratio_w, ratio_h))
        if ratio_str is None:
            # Check approximate match
            decimal = w / h
            if abs(decimal - 16 / 9) < 0.02:
                ratio_str = "16:9"
            elif abs(decimal - 4 / 3) < 0.02:
                ratio_str = "4:3"
            else:
                ratio_str = f"{ratio_w}:{ratio_h}"

        # Convert EMU to pixels (96 DPI, 914400 EMU per inch)
        px_w = round(int(w) / 914400 * 96)
        px_h = round(int(h) / 914400 * 96)

        return {
            'width_emu': int(w),
            'height_emu': int(h),
            'width_px': px_w,
            'height_px': px_h,
            'aspect_ratio': ratio_str,
        }

    def extract_color_palette(self, colors: dict[str, str]) -> dict[str, str]:
        """Map raw PPTX scheme colors to a semantic color palette.

        The semantic palette provides designer-friendly names that match
        common design-system tokens (primary, secondary, accent, etc.).
        """
        palette: dict[str, str] = {}

        # Background: dk2 if dark, else dk1
        for key in ('dk2', 'dk1'):
            if key in colors:
                r, g, b = hex_to_rgb(colors[key])
                if rgb_to_luminance(r, g, b) < 0.3:
                    palette['background'] = colors[key]
                    break
        if 'background' not in palette:
            palette['background'] = colors.get('lt1', '#FFFFFF')

        # Determine if the template is dark or light
        bg_lum = 1.0
        if 'background' in palette:
            r, g, b = hex_to_rgb(palette['background'])
            bg_lum = rgb_to_luminance(r, g, b)
        is_dark = bg_lum < 0.2

        # Text colours: choose high contrast against background
        if is_dark:
            palette['text_primary'] = colors.get('lt1', '#FFFFFF')
            palette['text_secondary'] = colors.get('lt2', '#B0B0B0')
        else:
            palette['text_primary'] = colors.get('dk1', '#000000')
            palette['text_secondary'] = colors.get('dk2', '#333333')

        # Primary = accent1, Secondary = accent2, Accent = accent3 or accent1
        palette['primary'] = colors.get('accent1', '#0078D4')
        palette['secondary'] = colors.get('accent2', palette['primary'])
        palette['accent'] = colors.get('accent3', palette['primary'])

        # Additional semantic tokens
        if 'accent4' in colors:
            palette['danger'] = colors['accent4']
        if 'accent5' in colors:
            palette['warning'] = colors['accent5']
        if 'accent6' in colors:
            palette['info'] = colors['accent6']
        if 'hlink' in colors:
            palette['link'] = colors['hlink']

        return palette

    def extract_header_region(self) -> dict[str, Any]:
        """Detect the header region of the slide master.

        Scans the top 15 % of the slide for shapes (text, lines, images)
        that appear on every slide, and returns them as structured
        ``header.elements``.
        """
        header_threshold = 15.0  # top 15 %
        elements: list[dict[str, Any]] = []
        min_top = 100.0
        max_bottom = 0.0

        for shape in self.master.shapes:
            top_pct = emu_to_percent(shape.top, SLIDE_HEIGHT_EMU)
            bottom_pct = emu_to_percent(shape.top + shape.height, SLIDE_HEIGHT_EMU)

            if top_pct > header_threshold:
                continue

            elem: dict[str, Any] = {
                'shape_name': shape.name,
                'position': {
                    'left_percent': emu_to_percent(shape.left, SLIDE_WIDTH_EMU),
                    'top_percent': top_pct,
                },
                'size': {
                    'width_percent': emu_to_percent(shape.width, SLIDE_WIDTH_EMU),
                    'height_percent': emu_to_percent(shape.height, SLIDE_HEIGHT_EMU),
                },
            }

            # Classify element type
            if shape.is_placeholder:
                ph = shape.placeholder_format
                elem['type'] = 'placeholder'
                elem['placeholder_type'] = str(ph.type) if ph.type else None
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                elem['type'] = 'image'
            elif hasattr(shape, 'text_frame'):
                text = ''
                try:
                    text = shape.text_frame.text.strip()
                except Exception:
                    pass
                if text:
                    elem['type'] = 'text'
                    elem['text'] = text
                else:
                    elem['type'] = 'shape'
            else:
                elem['type'] = 'shape'

            # Detect line / connector (divider)
            try:
                shape_type_name = str(shape.shape_type)
                if 'LINE' in shape_type_name or 'CONNECTOR' in shape_type_name:
                    elem['type'] = 'divider'
            except Exception:
                pass

            # Track vertical extent
            if top_pct < min_top:
                min_top = top_pct
            if bottom_pct > max_bottom:
                max_bottom = bottom_pct

            elements.append(elem)

        height_pct = round(max_bottom - min_top, 2) if elements else 0

        return {
            'position': 'top',
            'height_percent': height_pct,
            'elements': elements,
        }

    def extract_footer_region(self, footer_info: Optional[dict],
                               master_texts: list[dict]) -> dict[str, Any]:
        """Build a structured footer region from placeholders + master shapes.

        Similar to ``extract_header_region`` but for the bottom 15 % of the
        slide.  Also incorporates the existing footer/slide-number/date
        placeholders so everything is in one place.
        """
        footer_threshold = 85.0  # bottom 15 %
        elements: list[dict[str, Any]] = []
        min_top = 100.0
        max_bottom = 0.0

        for shape in self.master.shapes:
            top_pct = emu_to_percent(shape.top, SLIDE_HEIGHT_EMU)
            bottom_pct = emu_to_percent(shape.top + shape.height, SLIDE_HEIGHT_EMU)

            if top_pct < footer_threshold:
                continue

            elem: dict[str, Any] = {
                'shape_name': shape.name,
                'position': {
                    'left_percent': emu_to_percent(shape.left, SLIDE_WIDTH_EMU),
                    'top_percent': top_pct,
                },
                'size': {
                    'width_percent': emu_to_percent(shape.width, SLIDE_WIDTH_EMU),
                    'height_percent': emu_to_percent(shape.height, SLIDE_HEIGHT_EMU),
                },
            }

            if shape.is_placeholder:
                ph = shape.placeholder_format
                elem['type'] = 'placeholder'
                elem['placeholder_type'] = str(ph.type) if ph.type else None
                try:
                    if hasattr(shape, 'text_frame'):
                        elem['text'] = shape.text_frame.text.strip()
                except Exception:
                    pass
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                elem['type'] = 'image'
            elif hasattr(shape, 'text_frame'):
                text = ''
                try:
                    text = shape.text_frame.text.strip()
                except Exception:
                    pass
                if text:
                    elem['type'] = 'text'
                    elem['text'] = text
                else:
                    elem['type'] = 'shape'
            else:
                elem['type'] = 'shape'

            # Detect divider lines
            try:
                shape_type_name = str(shape.shape_type)
                if 'LINE' in shape_type_name or 'CONNECTOR' in shape_type_name:
                    elem['type'] = 'divider'
            except Exception:
                pass

            if top_pct < min_top:
                min_top = top_pct
            if bottom_pct > max_bottom:
                max_bottom = bottom_pct

            elements.append(elem)

        height_pct = round(max_bottom - min_top, 2) if elements else 0

        # Consolidated text (reuse existing logic)
        resolved_text = self._resolve_footer_text(footer_info, master_texts)

        return {
            'position': 'bottom',
            'height_percent': height_pct,
            'resolved_text': resolved_text,
            'elements': elements,
        }

    def _extract_shape_fill(self, shape) -> dict[str, Any]:
        """Extract fill information from a shape (solid, gradient, pattern).

        Returns a dict with ``fill_type`` and type-specific fields:
        - solid: ``fill_color``
        - gradient: ``gradient`` with ``angle`` and ``stops``
        - pattern/background: ``fill_type`` only
        """
        result: dict[str, Any] = {}
        try:
            fill = shape.fill
            if fill is None or fill.type is None:
                return result
            fill_type_str = str(fill.type)

            if 'SOLID' in fill_type_str:
                result['fill_type'] = 'solid'
                try:
                    fore = fill.fore_color
                    if fore and fore.rgb:
                        result['fill_color'] = f"#{fore.rgb}"
                except Exception:
                    pass

            elif 'GRADIENT' in fill_type_str:
                result['fill_type'] = 'gradient'
                grad_info: dict[str, Any] = {}
                try:
                    # Extract gradient stops from XML
                    a_ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
                    gs_list = shape._element.findall('.//' + a_ns + 'gs')
                    stops = []
                    for gs in gs_list:
                        pos = gs.get('pos')
                        pos_val = int(pos) / 1000 if pos else 0  # pos is in 1/1000ths of %
                        color = None
                        # Try srgbClr first, then schemeClr
                        srgb = gs.find('.//' + a_ns + 'srgbClr')
                        if srgb is not None:
                            color = f"#{srgb.get('val')}"
                        else:
                            scheme = gs.find('.//' + a_ns + 'schemeClr')
                            if scheme is not None:
                                scheme_name = scheme.get('val', '')
                                # Map scheme color to extracted theme colors
                                scheme_map = {
                                    'dk1': 'dk1', 'dk2': 'dk2',
                                    'lt1': 'lt1', 'lt2': 'lt2',
                                    'tx1': 'dk1', 'tx2': 'dk2',
                                    'bg1': 'lt1', 'bg2': 'lt2',
                                    'accent1': 'accent1', 'accent2': 'accent2',
                                    'accent3': 'accent3', 'accent4': 'accent4',
                                    'accent5': 'accent5', 'accent6': 'accent6',
                                    'hlink': 'hlink', 'folHlink': 'folHlink',
                                }
                                mapped = scheme_map.get(scheme_name)
                                if mapped and hasattr(self, '_theme_colors'):
                                    color = self._theme_colors.get(mapped)
                                elif mapped:
                                    color = f"scheme:{scheme_name}"
                                else:
                                    color = f"scheme:{scheme_name}"
                        if color:
                            stops.append({'pos': round(pos_val, 1), 'color': color})
                    if stops:
                        grad_info['stops'] = stops
                except Exception:
                    pass
                try:
                    lin = shape._element.find(
                        './/' + '{http://schemas.openxmlformats.org/drawingml/2006/main}lin')
                    if lin is not None:
                        ang = lin.get('ang')
                        if ang:
                            grad_info['angle'] = int(ang) // 60000  # EMU to degrees
                except Exception:
                    pass
                if grad_info:
                    result['gradient'] = grad_info

            elif 'PATTERN' in fill_type_str:
                result['fill_type'] = 'pattern'
            elif 'BACKGROUND' in fill_type_str:
                result['fill_type'] = 'background'
        except Exception:
            pass
        return result

    def _extract_preset_geometry(self, shape) -> Optional[str]:
        """Extract preset geometry name from a shape (e.g., chevron, arc, ellipse)."""
        try:
            prst = shape._element.find(
                './/' + '{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
            if prst is not None:
                return prst.get('prst')
        except Exception:
            pass
        return None

    def _build_decorative_elem(self, shape, source: str = 'master') -> dict[str, Any]:
        """Build a decorative element dict from a shape."""
        top_pct = emu_to_percent(shape.top, SLIDE_HEIGHT_EMU)
        left_pct = emu_to_percent(shape.left, SLIDE_WIDTH_EMU)
        w_pct = emu_to_percent(shape.width, SLIDE_WIDTH_EMU)
        h_pct = emu_to_percent(shape.height, SLIDE_HEIGHT_EMU)

        shape_type_str = ''
        try:
            shape_type_str = str(shape.shape_type)
        except Exception:
            pass

        elem: dict[str, Any] = {
            'shape_name': shape.name,
            'shape_type': shape_type_str,
            'position': {
                'left_percent': left_pct,
                'top_percent': top_pct,
            },
            'size': {
                'width_percent': w_pct,
                'height_percent': h_pct,
            },
            'source': source,
        }

        # Preset geometry
        geo = self._extract_preset_geometry(shape)
        if geo:
            elem['preset_geometry'] = geo

        # Fill info
        fill_info = self._extract_shape_fill(shape)
        elem.update(fill_info)

        # Rotation
        try:
            if shape.rotation and shape.rotation != 0:
                elem['rotation'] = shape.rotation
        except Exception:
            pass

        return elem

    def _is_decorative_candidate(self, shape) -> bool:
        """Check if a shape qualifies as a decorative element."""
        if shape.is_placeholder:
            return False

        # Skip shapes whose names indicate structural (non-decorative) roles
        name_lower = shape.name.lower()
        if any(kw in name_lower for kw in ('placeholder', 'slide number', 'footer', 'date')):
            return False

        # Skip practically invisible shapes (< 0.5% in either dimension)
        # but allow lines/connectors which are 1-dimensional
        w_pct = emu_to_percent(shape.width, SLIDE_WIDTH_EMU)
        h_pct = emu_to_percent(shape.height, SLIDE_HEIGHT_EMU)
        shape_type_str = str(shape.shape_type) if shape.shape_type else ''
        is_line = 'LINE' in shape_type_str or 'CONNECTOR' in shape_type_str
        if not is_line and (w_pct < 0.5 or h_pct < 0.5):
            return False

        # Skip shapes entirely off-screen
        left_pct = emu_to_percent(shape.left, SLIDE_WIDTH_EMU)
        top_pct = emu_to_percent(shape.top, SLIDE_HEIGHT_EMU)
        if left_pct + w_pct < 0 or top_pct + h_pct < 0:
            return False

        logo_threshold = SLIDE_WIDTH_EMU * 0.20

        # Skip small pictures (logos) and large pictures (backgrounds)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if shape.width < logo_threshold and shape.height < logo_threshold:
                return False
            if (shape.width > SLIDE_WIDTH_EMU * 0.8
                    and shape.height > SLIDE_HEIGHT_EMU * 0.8):
                return False

        # Allow shapes with short text (≤3 chars, e.g. ">" chevron labels)
        # but skip shapes with meaningful text
        if hasattr(shape, 'text_frame'):
            try:
                text = shape.text_frame.text.strip()
                if len(text) > 3:
                    return False
            except Exception:
                pass

        return True

    def extract_decorative_elements(self) -> list[dict[str, Any]]:
        """Detect decorative shapes on the slide master and common layouts.

        Decorative elements are non-placeholder, non-text, non-logo shapes
        that serve a purely visual purpose (chevrons, dot patterns,
        semicircles, colored bars, gradient rectangles, etc.).

        Scans both the slide master and all layouts.  Master-level shapes
        are always included.  Layout-level shapes are included when the
        same shape name appears in multiple layouts (≥ 3 or ≥ 5% of
        layouts, whichever is larger).
        """
        decorative: list[dict[str, Any]] = []

        # --- 1. Master shapes ---
        for shape in self.master.shapes:
            if self._is_decorative_candidate(shape):
                decorative.append(self._build_decorative_elem(shape, source='master'))

        # --- 2. Layout common shapes ---
        # Track shapes by name+approximate position across all layouts
        layout_shapes: dict[str, list] = {}  # name -> list of (layout_idx, shape)
        layouts = list(self.master.slide_layouts)
        total_layouts = len(layouts)

        for i, layout in enumerate(layouts):
            for shape in layout.shapes:
                if not self._is_decorative_candidate(shape):
                    continue
                key = shape.name
                if key not in layout_shapes:
                    layout_shapes[key] = []
                layout_shapes[key].append((i, shape))

        threshold = max(3, int(total_layouts * 0.05))
        seen_names: set[str] = set()
        for name, entries in layout_shapes.items():
            if len(entries) >= threshold and name not in seen_names:
                seen_names.add(name)
                # Use the first occurrence as representative
                _, rep_shape = entries[0]
                elem = self._build_decorative_elem(rep_shape, source='layout_common')
                elem['layout_count'] = len(entries)
                elem['layout_total'] = total_layouts
                decorative.append(elem)

        return decorative

    def extract_layout_decorative_patterns(self) -> list[dict[str, Any]]:
        """Summarise decorative patterns found across all layouts.

        Groups shapes by category (gradient_panel, circle_icon, divider,
        rounded_accent, decorative_shape) and records representative
        parameters for each category.
        """
        categories: dict[str, list[dict[str, Any]]] = {
            'gradient_panel': [],
            'circle_icon': [],
            'divider': [],
            'rounded_accent': [],
            'decorative_shape': [],
        }

        for layout in self.master.slide_layouts:
            for shape in layout.shapes:
                if not self._is_decorative_candidate(shape):
                    continue

                geo = self._extract_preset_geometry(shape)
                fill = self._extract_shape_fill(shape)
                w_pct = emu_to_percent(shape.width, SLIDE_WIDTH_EMU)
                h_pct = emu_to_percent(shape.height, SLIDE_HEIGHT_EMU)
                shape_type_str = str(shape.shape_type) if shape.shape_type else ''

                entry = {
                    'layout': layout.name,
                    'shape_name': shape.name,
                    'preset_geometry': geo,
                    'size': {'width_percent': w_pct, 'height_percent': h_pct},
                }
                entry.update(fill)

                # Classify
                if fill.get('fill_type') == 'gradient' and (w_pct > 20 or h_pct > 20):
                    categories['gradient_panel'].append(entry)
                elif geo == 'ellipse':
                    categories['circle_icon'].append(entry)
                elif 'LINE' in shape_type_str or geo == 'line':
                    categories['divider'].append(entry)
                elif geo == 'roundRect':
                    categories['rounded_accent'].append(entry)
                else:
                    categories['decorative_shape'].append(entry)

        patterns = []
        for cat, items in categories.items():
            if not items:
                continue
            patterns.append({
                'category': cat,
                'count': len(items),
                'examples': items[:3],  # Up to 3 representative examples
            })

        return patterns

    def build_slide_master(self, manifest: dict[str, Any],
                           design_source: Optional[dict[str, Any]] = None) -> dict[str, Any]:
        """Assemble the ``slide_master`` section from extracted data.

        This is the structured representation of the common frame that
        applies to every slide (as opposed to per-layout content areas).

        Args:
            manifest: The partially-built theme manifest.
            design_source: Optional external design source reference
                (Figma, Google Slides URL with metadata).
        """
        colors = manifest.get('colors', {})
        footer_info = manifest.get('footer')
        master_texts = manifest.get('master_texts', [])

        result = {
            'slide_size': self.extract_slide_size(),
            'color_palette': self.extract_color_palette(colors),
            'header': self.extract_header_region(),
            'footer': self.extract_footer_region(footer_info, master_texts),
            'decorative_elements': self.extract_decorative_elements(),
            'layout_decorative_patterns': self.extract_layout_decorative_patterns(),
        }

        if design_source:
            result['design_source'] = design_source

        return result

    def extract_layout_details(self) -> list[dict[str, Any]]:
        """Extract details from all slide master layouts.

        Includes background type, placeholders, and non-placeholder text
        shapes for each layout. Useful for referencing the original PPTX
        structure when building HTML slides.
        """
        layouts = []
        for i, layout in enumerate(self.master.slide_layouts):
            layout_info = {
                'index': i,
                'name': layout.name,
                'background': self._extract_background_from_element(layout),
                'placeholders': [],
                'texts': [],
            }

            for shape in layout.shapes:
                if shape.is_placeholder:
                    ph = shape.placeholder_format
                    ph_info = {
                        'idx': ph.idx,
                        'type': str(ph.type) if ph.type else None,
                        'name': shape.name,
                    }
                    try:
                        if hasattr(shape, 'text_frame'):
                            text = shape.text_frame.text.strip()
                            if text:
                                ph_info['text'] = text
                    except Exception:
                        pass
                    layout_info['placeholders'].append(ph_info)
                elif hasattr(shape, 'text_frame'):
                    try:
                        text = shape.text_frame.text.strip()
                        if text:
                            layout_info['texts'].append({
                                'text': text,
                                'shape_name': shape.name,
                            })
                    except Exception:
                        pass

            layouts.append(layout_info)

        return layouts

    def _resolve_footer_text(self, footer_info: Optional[dict],
                              master_texts: list[dict]) -> Optional[str]:
        """Resolve consolidated footer text from placeholders and master texts.

        Combines footer placeholder text with footer-area master texts,
        deduplicating identical messages. Returns a single string or None.
        """
        parts = []
        seen = set()

        # Footer placeholder text
        if footer_info and footer_info.get('text'):
            text = footer_info['text'].strip()
            if text and text not in seen:
                parts.append(text)
                seen.add(text)

        # Master texts in footer area (bottom 15% of slide)
        for mt in master_texts:
            if mt.get('is_footer_area'):
                text = mt['text'].strip()
                if text and text not in seen:
                    parts.append(text)
                    seen.add(text)

        return ' | '.join(parts) if parts else None

    def list_masters(self) -> list[dict[str, Any]]:
        """List all slide masters with layout counts."""
        masters = []
        for i, master in enumerate(self.prs.slide_masters):
            masters.append({
                'index': i,
                'name': master.name if hasattr(master, 'name') else f"Master {i}",
                'layout_count': len(master.slide_layouts)
            })
        return masters

    def list_layouts(self) -> list[dict[str, Any]]:
        """List all layouts in selected master with background types."""
        layouts = []
        for i, layout in enumerate(self.master.slide_layouts):
            bg_info = self._extract_background_from_element(layout)
            layouts.append({
                'index': i,
                'name': layout.name,
                'background_type': bg_info['type']
            })
        return layouts

    def extract_all(self, output_dir: Path,
                    design_source: Optional[dict[str, Any]] = None) -> dict[str, Any]:
        """
        Extract all theme information.

        Args:
            output_dir: Directory to save extracted assets
            design_source: Optional external design-guide reference
                (Figma URL, Google Slides URL, etc.).

        Returns:
            Complete theme manifest dictionary.
        """
        output_dir.mkdir(parents=True, exist_ok=True)

        manifest = {
            'source_file': str(self.pptx_path.name),
            'theme_name': self.theme_name,
            'master_index': self.master_index,
            'master_name': self.master.name if hasattr(self.master, 'name') else None,
            'slide_count': len(self.prs.slides),
            'layout_count': len(self.master.slide_layouts),
            'colors': self.extract_colors(),
            'fonts': self.extract_fonts(),
            'backgrounds': self.extract_backgrounds(output_dir),
            'logos': self.extract_logos(output_dir),
            'footer': self.extract_footer(),
            'slide_number': self.extract_slide_number(),
            'date': self.extract_date(),
            'master_texts': self.extract_master_texts(),
            'layout_details': self.extract_layout_details(),
        }

        # Resolved footer text: deduplicated from placeholder + master text shapes
        manifest['footer_text'] = self._resolve_footer_text(
            manifest['footer'], manifest['master_texts']
        )

        # Slide-master template: common frame across all slides
        manifest['slide_master'] = self.build_slide_master(manifest, design_source)

        return manifest


class CSSGenerator:
    """Generates CSS overrides from extracted theme data."""

    # System fonts that don't need web font loading
    SYSTEM_FONTS = {
        'Arial', 'Helvetica', 'Times New Roman', 'Times', 'Courier New',
        'Courier', 'Verdana', 'Georgia', 'Palatino', 'Garamond', 'Bookman',
        'Comic Sans MS', 'Trebuchet MS', 'Arial Black', 'Impact', 'Calibri',
        'Calibri Light', 'Cambria', 'Consolas', 'Segoe UI'
    }

    def __init__(self, manifest: dict[str, Any]):
        """
        Initialize CSS generator with theme manifest.

        Args:
            manifest: Theme manifest dictionary from ThemeExtractor
        """
        self.manifest = manifest
        self.colors = manifest.get('colors', {})
        self.fonts = manifest.get('fonts', {})
        self.logos = manifest.get('logos', [])
        self.footer = manifest.get('footer')
        self.slide_number = manifest.get('slide_number')
        self.backgrounds = manifest.get('backgrounds', {})
        self.master_texts = manifest.get('master_texts', [])

    def generate(self) -> str:
        """Generate complete CSS override file."""
        sections = [
            self._generate_header(),
            self._generate_color_variables(),
            self._generate_font_variables(),
            self._generate_logo_styles(),
            self._generate_footer_styles(),
            self._generate_slide_number_styles(),
            self._generate_background_styles(),
        ]

        return '\n\n'.join(filter(None, sections))

    def _generate_header(self) -> str:
        """Generate CSS file header comment."""
        return f"""/*
 * Theme Override CSS
 * Generated from: {self.manifest.get('source_file', 'Unknown')}
 * Theme: {self.manifest.get('theme_name', 'Unknown')}
 * Generated by extract_pptx_theme.py
 */"""

    def _generate_color_variables(self) -> str:
        """Generate CSS color variables from PPTX color scheme."""
        if not self.colors:
            return ""

        lines = ["/* Color Variables */", ":root {"]

        # Map PPTX colors to CSS variables
        color_mapping = {
            'accent1': '--accent',
            'accent2': '--accent-light',
            'accent3': '--green',
            'accent4': '--red',
            'accent5': '--orange',
            'accent6': '--yellow',
            'hlink': '--cyan',
        }

        # Check if dk2 is dark enough for background
        dk2 = self.colors.get('dk2', '')
        if dk2:
            r, g, b = hex_to_rgb(dk2)
            luminance = rgb_to_luminance(r, g, b)
            if luminance < 0.2:
                lines.append(f"  --bg-primary: {dk2};")
            else:
                lines.append("  /* dk2 too light for dark theme, keeping default */")
                lines.append("  --bg-primary: #0f1117;")

        # Map accent colors
        for pptx_name, css_var in color_mapping.items():
            if pptx_name in self.colors:
                lines.append(f"  {css_var}: {self.colors[pptx_name]};")

        # Generate accent glow from accent1
        accent1 = self.colors.get('accent1', '#41B3FF')
        r, g, b = hex_to_rgb(accent1)
        lines.append(f"  --accent-glow: rgba({r}, {g}, {b}, 0.3);")

        # Keep text colors light for dark background
        lines.append("")
        lines.append("  /* Text colors (kept light for dark background) */")
        lines.append("  --text-primary: #ffffff;")
        lines.append("  --text-secondary: #b0b0b0;")

        # Include original PPTX colors as reference
        lines.append("")
        lines.append("  /* Original PPTX theme colors (reference) */")
        for name, value in self.colors.items():
            lines.append(f"  --pptx-{name}: {value};")

        lines.append("}")

        return '\n'.join(lines)

    def _generate_font_variables(self) -> str:
        """Generate CSS font variables."""
        if not self.fonts:
            return ""

        lines = ["/* Font Variables */"]

        heading_font = self.fonts.get('heading', 'Calibri Light')
        body_font = self.fonts.get('body', 'Calibri')

        # Check if fonts need web font loading
        needs_web_font = []
        if heading_font not in self.SYSTEM_FONTS:
            needs_web_font.append(heading_font)
        if body_font not in self.SYSTEM_FONTS and body_font != heading_font:
            needs_web_font.append(body_font)

        if needs_web_font:
            lines.append(f"/* NOTE: The following fonts may need web font loading: {', '.join(needs_web_font)} */")
            lines.append("/* Consider adding Google Fonts or @font-face declarations */")

        lines.append(":root {")
        lines.append(f"  --font-heading: '{heading_font}', system-ui, sans-serif;")
        lines.append(f"  --font-body: '{body_font}', system-ui, sans-serif;")
        lines.append("}")

        return '\n'.join(lines)

    def _generate_logo_styles(self) -> str:
        """Generate CSS for logo positioning."""
        if not self.logos:
            return ""

        lines = ["/* Logo Styles */"]

        for i, logo in enumerate(self.logos):
            selector = ".slide-logo" if i == 0 else f".slide-logo-{i + 1}"
            pos = logo.get('position', {})
            size = logo.get('size', {})

            lines.append(f"{selector} {{")
            lines.append("  position: absolute;")
            lines.append(f"  left: {pos.get('left_percent', 5)}%;")
            lines.append(f"  top: {pos.get('top_percent', 90)}%;")
            lines.append(f"  width: {size.get('width_percent', 3)}%;")
            lines.append(f"  height: auto;")
            lines.append(f"  background-image: url('./images/{logo.get('filename', 'logo.png')}');")
            lines.append("  background-size: contain;")
            lines.append("  background-repeat: no-repeat;")
            lines.append("}")

            # Add nearby text as comment
            nearby_text = logo.get('nearby_text')
            if nearby_text:
                lines.append(f"/* Nearby text: {nearby_text} */")

        return '\n'.join(lines)

    def _generate_footer_styles(self) -> str:
        """Generate CSS for footer positioning.

        Footer text is set via SlideFramework({ footer: '...' }) in JS,
        not via CSS content property. CSS only handles positioning.
        """
        if not self.footer:
            return ""

        lines = ["/* Footer Styles */"]
        pos = self.footer.get('position', {})
        text = self.footer.get('text', '')

        if text:
            escaped = text.replace("'", "\\'")
            lines.append(f"/* Footer text from PPTX: \"{text}\" */")
            lines.append(f"/* Use in SlideFramework: new SlideFramework({{ footer: '{escaped}' }}) */")

        # Also check master_texts for additional footer-area text
        master_texts = self.manifest.get('master_texts', [])
        footer_area_texts = [mt['text'] for mt in master_texts if mt.get('is_footer_area')]
        if footer_area_texts:
            lines.append(f"/* Additional master footer text: {footer_area_texts} */")

        footer_text = self.manifest.get('footer_text')
        if footer_text:
            escaped_ft = footer_text.replace("'", "\\'")
            lines.append(f"/* Resolved footer_text: '{escaped_ft}' */")

        lines.append(".slide-footer {")
        lines.append("  position: absolute;")
        lines.append(f"  left: {pos.get('left_percent', 5)}%;")
        lines.append(f"  bottom: {100 - pos.get('top_percent', 95)}%;")
        lines.append("  font-size: 0.7rem;")
        lines.append("  color: var(--text-secondary);")
        lines.append("}")

        return '\n'.join(lines)

    def _generate_slide_number_styles(self) -> str:
        """Generate CSS for slide number."""
        if not self.slide_number:
            return ""

        lines = ["/* Slide Number Styles */"]
        pos = self.slide_number.get('position', {})
        text_format = self.slide_number.get('text', '')

        lines.append(".slide-counter {")
        lines.append("  position: absolute;")
        lines.append(f"  right: {100 - pos.get('left_percent', 95) - (self.slide_number.get('size', {}).get('width_percent', 5))}%;")
        lines.append(f"  bottom: {100 - pos.get('top_percent', 95)}%;")
        lines.append("  font-size: 0.7rem;")
        lines.append("  color: var(--text-secondary);")
        lines.append("}")

        if text_format:
            lines.append(f"/* Original format: {text_format} */")

        return '\n'.join(lines)

    def _generate_background_styles(self) -> str:
        """Generate CSS for background images."""
        master_bg = self.backgrounds.get('master', {})

        if master_bg.get('type') != 'picture':
            return ""

        lines = ["/* Background Image Styles */"]
        lines.append(".slide {")
        lines.append("  position: relative;")
        lines.append("}")
        lines.append("")
        lines.append(".slide::before {")
        lines.append("  content: '';")
        lines.append("  position: absolute;")
        lines.append("  top: 0;")
        lines.append("  left: 0;")
        lines.append("  right: 0;")
        lines.append("  bottom: 0;")
        lines.append("  background-image: url('./images/background.png');")
        lines.append("  background-size: cover;")
        lines.append("  background-position: center;")
        lines.append("  /* Dark overlay for readability */")
        lines.append("  filter: brightness(0.3);")
        lines.append("  z-index: -1;")
        lines.append("}")
        lines.append("")
        lines.append("/* Alternative: overlay approach */")
        lines.append("/*")
        lines.append(".slide-background-overlay {")
        lines.append("  position: absolute;")
        lines.append("  top: 0;")
        lines.append("  left: 0;")
        lines.append("  right: 0;")
        lines.append("  bottom: 0;")
        lines.append("  background: rgba(15, 17, 23, 0.85);")
        lines.append("  z-index: 0;")
        lines.append("}")
        lines.append("*/")

        return '\n'.join(lines)


class DesignGuideGenerator:
    """Generates a design.md design-system document from theme manifest.

    Inspired by Figma's design-token philosophy: capture the *rules* behind
    visual choices so that every new slide is automatically consistent.
    The output is a Markdown file with YAML code blocks for machine-readable
    tokens plus prose descriptions for human designers and AI agents.
    """

    def __init__(self, manifest: dict[str, Any]):
        self.m = manifest
        self.sm = manifest.get('slide_master', {})
        self.colors = manifest.get('colors', {})
        self.palette = self.sm.get('color_palette', {})
        self.fonts = manifest.get('fonts', {})
        self.logos = manifest.get('logos', [])
        self.header = self.sm.get('header', {})
        self.footer_region = self.sm.get('footer', {})
        self.decorative = self.sm.get('decorative_elements', [])
        self.layout_patterns = self.sm.get('layout_decorative_patterns', [])
        self.slide_size = self.sm.get('slide_size', {})
        self.backgrounds = manifest.get('backgrounds', {})
        self.design_source = self.sm.get('design_source')

    def generate(self) -> str:
        sections = [
            self._section_header(),
            self._section_slide_canvas(),
            self._section_color_system(),
            self._section_typography(),
            self._section_iconography(),
            self._section_spacing_grid(),
            self._section_shapes_and_corners(),
            self._section_decorative_patterns(),
            self._section_header_footer(),
            self._section_logo_branding(),
            self._section_backgrounds(),
            self._section_motion(),
            self._section_design_source(),
            self._section_usage_rules(),
        ]
        return '\n\n'.join(filter(None, sections)) + '\n'

    # -- Individual sections ------------------------------------------------

    def _section_header(self) -> str:
        theme_name = self.m.get('theme_name', 'Untitled')
        source = self.m.get('source_file', 'unknown.pptx')
        return f"""# Design System — {theme_name}

> Auto-generated from `{source}` by `extract_pptx_theme.py`.
> This document defines the visual language that must be followed when
> creating new slides.  Treat every token below as a **constraint**, not
> a suggestion.  Consistency > creativity."""

    def _section_slide_canvas(self) -> str:
        ss = self.slide_size
        if not ss:
            return ''
        w = ss.get('width_px', 1280)
        h = ss.get('height_px', 720)
        return f"""## 1. Slide Canvas

| Property | Value |
|----------|-------|
| Aspect ratio | **{ss.get('aspect_ratio', '16:9')}** |
| Resolution | {w} × {h} px |
| EMU | {ss.get('width_emu', '')} × {ss.get('height_emu', '')} |
| Remarp frontmatter | `size: {w}x{h}` |

All content must be designed for this aspect ratio.  Do not stretch or
crop to fit a different ratio.

Add this to your `.remarp.md` frontmatter:
```yaml
size: {w}x{h}
```"""

    def _section_color_system(self) -> str:
        if not self.palette:
            return ''

        lines = [
            '## 2. Color System',
            '',
            'Semantic color tokens mapped from the PPTX theme scheme.',
            'Use semantic names (not raw hex) so the palette can be',
            'swapped without touching individual slides.',
            '',
            '```yaml',
            'color_tokens:',
        ]

        role_desc = {
            'primary': 'Primary brand / accent (CTA, links, highlights)',
            'secondary': 'Secondary accent (gradients, hover, subtle emphasis)',
            'accent': 'Tertiary accent (success, positive indicators)',
            'background': 'Slide background',
            'text_primary': 'Headings and body text',
            'text_secondary': 'Captions, labels, muted text',
            'danger': 'Error / destructive actions',
            'warning': 'Caution indicators',
            'info': 'Informational highlights',
            'link': 'Hyperlinks and interactive text',
        }

        for key in ('background', 'text_primary', 'text_secondary',
                     'primary', 'secondary', 'accent',
                     'danger', 'warning', 'info', 'link'):
            val = self.palette.get(key)
            if val:
                desc = role_desc.get(key, '')
                lines.append(f'  {key}: "{val}"  # {desc}')

        lines.append('```')
        lines.append('')

        # Contrast guidance
        bg = self.palette.get('background', '#000000')
        r, g, b = hex_to_rgb(bg)
        lum = rgb_to_luminance(r, g, b)
        if lum < 0.2:
            lines.append('**Theme type**: Dark.  Text must be light (#fff / #ccc).')
            lines.append('Avoid pure white on pure black — use the token values above.')
        else:
            lines.append('**Theme type**: Light.  Text must be dark (#000 / #333).')
            lines.append('Ensure sufficient contrast (WCAG AA ≥ 4.5:1 for body text).')

        # Raw PPTX scheme reference
        if self.colors:
            lines.append('')
            lines.append('<details><summary>Raw PPTX scheme colors (reference)</summary>')
            lines.append('')
            lines.append('| Slot | Hex |')
            lines.append('|------|-----|')
            for k, v in self.colors.items():
                lines.append(f'| {k} | `{v}` |')
            lines.append('')
            lines.append('</details>')

        return '\n'.join(lines)

    def _section_typography(self) -> str:
        if not self.fonts:
            return ''

        heading = self.fonts.get('heading', 'system-ui')
        body = self.fonts.get('body', 'system-ui')

        return f"""## 3. Typography

```yaml
typography:
  heading:
    family: "{heading}"
    weights: [400, 700]       # regular + bold
    usage: "Slide titles, section headers, callout text"
  body:
    family: "{body}"
    weights: [400, 600]       # regular + semibold
    usage: "Body text, bullets, labels, speaker notes"
```

### Type Scale (recommended)

| Level | Size | Weight | Use |
|-------|------|--------|-----|
| H1 | 2.5 rem | 700 | Cover slide title |
| H2 | 1.8 rem | 700 | Slide title |
| H3 | 1.3 rem | 600 | Sub-heading |
| Body | 1.0 rem | 400 | Bullet text, paragraphs |
| Caption | 0.8 rem | 400 | Footer, labels, annotations |
| Code | 0.9 rem (mono) | 400 | Code blocks |

> If `{heading}` or `{body}` is not a system font, add a `@import` or
> `@font-face` declaration.  Never fall back to a visually different
> typeface silently."""

    def _section_iconography(self) -> str:
        # Infer icon style from decorative elements and template characteristics
        has_rounded = any('round' in d.get('shape_name', '').lower()
                         for d in self.decorative)
        has_sharp = any('rect' in d.get('shape_name', '').lower()
                       or 'sharp' in d.get('shape_name', '').lower()
                       for d in self.decorative)

        # Default recommendation based on common patterns
        style = 'rounded'
        if has_sharp and not has_rounded:
            style = 'sharp'

        return f"""## 4. Iconography

```yaml
icons:
  style: "{style}"           # rounded | sharp | outlined | filled
  corner_radius: 4px         # icon container rounding
  size_default: 48px         # standard icon size on content slides
  size_small: 24px           # inline / label icons
  size_large: 64px           # hero / feature highlight
  color: "inherit"           # icons inherit text color by default
  accent_color: "primary"    # highlighted icons use the primary token
```

### Rules

- **Consistency**: all icons on a single slide must use the same style
  (`{style}`).  Do not mix rounded and sharp icons.
- **AWS service icons**: use the official SVG from
  `skills/reactive-presentation/assets/aws-icons/`.  Keep the original
  multi-color fill — do not monochrome AWS icons.
- **Generic icons**: prefer a single icon library (e.g. Lucide, Phosphor,
  Material Symbols) across the entire presentation.
- **Sizing**: icons next to text should be vertically centered and match
  the text line height.
- **Padding**: maintain ≥ 8 px clear space around every icon."""

    def _section_spacing_grid(self) -> str:
        ar = self.slide_size.get('aspect_ratio', '16:9')
        # Common safe-area margins
        return f"""## 5. Spacing & Grid

```yaml
layout:
  aspect_ratio: "{ar}"
  safe_area:
    top: 10%                  # below header region
    bottom: 10%               # above footer region
    left: 5%
    right: 5%
  grid:
    columns: 12
    gutter: 16px
    margin: 40px              # outer margin
spacing_scale:               # 4-px base unit
  xs: 4px
  sm: 8px
  md: 16px
  lg: 24px
  xl: 32px
  xxl: 48px
```

### Rules

- Content must **never** overlap the header or footer regions.
- Use multiples of the 4 px base unit for all padding and margins.
- Two-column layouts: use 50/50 split with `md` gutter.
- Three-column layouts: use 33/33/33 split with `md` gutter.
- Maintain consistent vertical rhythm — every element should snap
  to the spacing scale."""

    def _section_shapes_and_corners(self) -> str:
        # Infer corner radius from decorative elements
        return """## 6. Shapes & Corners

```yaml
shapes:
  border_radius:
    none: 0px                 # tables, code blocks
    sm: 4px                   # tags, badges, small chips
    md: 8px                   # cards, content boxes, tooltips
    lg: 16px                  # hero cards, feature panels
    full: 9999px              # pills, avatar circles
  border:
    width: 1px
    color: "text_secondary"   # use semantic token
    style: "solid"
  shadow:
    none: "none"
    sm: "0 1px 2px rgba(0,0,0,0.1)"
    md: "0 4px 12px rgba(0,0,0,0.15)"
    lg: "0 8px 24px rgba(0,0,0,0.2)"
```

### Rules

- **Cards and containers**: use `md` radius (8 px).
- **Buttons and badges**: use `sm` radius (4 px) or `full` for pills.
- **Images inside cards**: clip to parent's border-radius.
- **Do not mix** rounded and sharp containers on the same slide.
- **Elevation (shadow)**: use sparingly; on dark backgrounds prefer
  a subtle border or glow over a drop shadow."""

    def _section_decorative_patterns(self) -> str:
        lines = ['## 7. Decorative Patterns']

        if self.decorative:
            lines += [
                '',
                'The following decorative shapes were detected.',
                'Master-level shapes **must be replicated** on every slide.',
                '',
                '```yaml',
                'decorative_elements:',
            ]

            for i, d in enumerate(self.decorative):
                pos = d.get('position', {})
                size = d.get('size', {})
                lines.append(f'  - name: "{d.get("shape_name", f"element_{i}")}"')
                lines.append(f'    source: "{d.get("source", "master")}"')
                if d.get('preset_geometry'):
                    lines.append(f'    geometry: "{d["preset_geometry"]}"')
                lines.append(f'    type: "{d.get("shape_type", "unknown")}"')
                lines.append(f'    position: {{ left: {pos.get("left_percent", 0)}%, top: {pos.get("top_percent", 0)}% }}')
                lines.append(f'    size: {{ width: {size.get("width_percent", 0)}%, height: {size.get("height_percent", 0)}% }}')
                if d.get('fill_type'):
                    lines.append(f'    fill_type: "{d["fill_type"]}"')
                if d.get('fill_color'):
                    lines.append(f'    fill_color: "{d["fill_color"]}"')
                if d.get('gradient'):
                    grad = d['gradient']
                    lines.append(f'    gradient:')
                    if 'angle' in grad:
                        lines.append(f'      angle: {grad["angle"]}')
                    if 'stops' in grad:
                        lines.append(f'      stops: {grad["stops"]}')
                if d.get('rotation'):
                    lines.append(f'    rotation: {d["rotation"]}deg')
                if d.get('layout_count'):
                    lines.append(f'    layout_frequency: {d["layout_count"]}/{d["layout_total"]}')

            lines.append('```')
            lines.append('')
            lines.append('> Master-level shapes form the slide "frame".  Layout-common')
            lines.append('> shapes indicate prevalent design vocabulary in this template.')
        else:
            lines += [
                '',
                'No decorative elements detected on the slide master.',
            ]

        # Layout decorative patterns summary
        if self.layout_patterns:
            lines += [
                '',
                '### Layout Decorative Vocabulary',
                '',
                'The following decorative pattern categories were found across',
                'slide layouts.  Use these as the design vocabulary when adding',
                'visual elements to slides.',
                '',
            ]
            for pat in self.layout_patterns:
                cat = pat['category']
                count = pat['count']
                lines.append(f'**{cat}** ({count} instances)')
                lines.append('')
                for ex in pat.get('examples', []):
                    geo = ex.get('preset_geometry', 'rect')
                    size = ex.get('size', {})
                    fill = ex.get('fill_type', '')
                    layout = ex.get('layout', '')
                    lines.append(f'  - `{geo}` {size.get("width_percent", 0)}%×{size.get("height_percent", 0)}% fill={fill} (from "{layout}")')
                    if ex.get('gradient'):
                        grad = ex['gradient']
                        stops = grad.get('stops', [])
                        if stops:
                            colors = ' → '.join(s.get('color', '') for s in stops)
                            lines.append(f'    gradient: {colors}')
                lines.append('')
        elif not self.decorative:
            lines += [
                '',
                'Keep slides clean — avoid adding ornamental shapes unless the',
                'design guide is updated.',
            ]

        return '\n'.join(lines)

    def _section_header_footer(self) -> str:
        lines = ['## 8. Header & Footer Regions']

        # Header
        h = self.header
        if h and h.get('elements'):
            lines.append('')
            lines.append(f'### Header (top {h.get("height_percent", 0)}%)')
            lines.append('')
            lines.append('| Element | Type | Position (L%, T%) | Size (W%, H%) |')
            lines.append('|---------|------|--------------------|----------------|')
            for el in h['elements']:
                p = el.get('position', {})
                s = el.get('size', {})
                text = f' — "{el["text"]}"' if el.get('text') else ''
                lines.append(f'| {el.get("shape_name", "")} | {el.get("type", "")}{text} | {p.get("left_percent", 0)}, {p.get("top_percent", 0)} | {s.get("width_percent", 0)}, {s.get("height_percent", 0)} |')
        else:
            lines.append('\nNo header region elements detected.')

        # Footer
        f = self.footer_region
        if f and f.get('elements'):
            lines.append('')
            lines.append(f'### Footer (bottom, height {f.get("height_percent", 0)}%)')
            lines.append('')
            if f.get('resolved_text'):
                lines.append(f'**Footer text**: `{f["resolved_text"]}`')
                lines.append('')
            lines.append('| Element | Type | Position (L%, T%) | Size (W%, H%) |')
            lines.append('|---------|------|--------------------|----------------|')
            for el in f['elements']:
                p = el.get('position', {})
                s = el.get('size', {})
                text = f' — "{el["text"]}"' if el.get('text') else ''
                lines.append(f'| {el.get("shape_name", "")} | {el.get("type", "")}{text} | {p.get("left_percent", 0)}, {p.get("top_percent", 0)} | {s.get("width_percent", 0)}, {s.get("height_percent", 0)} |')
        else:
            lines.append('\nNo footer region elements detected.')

        lines.append('')
        lines.append('> Header and footer are **fixed regions**.  Slide content must')
        lines.append('> not intrude into these areas.')

        return '\n'.join(lines)

    def _section_logo_branding(self) -> str:
        if not self.logos:
            return '## 9. Logo & Branding\n\nNo logos detected on the slide master.'

        lines = ['## 9. Logo & Branding', '']

        for i, logo in enumerate(self.logos):
            pos = logo.get('position', {})
            size = logo.get('size', {})
            lines.append(f'### Logo {i + 1}: `{logo.get("filename", "logo.png")}`')
            lines.append('')
            lines.append(f'- **Position**: left {pos.get("left_percent", 0)}%, top {pos.get("top_percent", 0)}%')
            lines.append(f'- **Size**: {size.get("width_percent", 0)}% × {size.get("height_percent", 0)}%')
            nearby = logo.get('nearby_text')
            if nearby:
                lines.append(f'- **Nearby text**: "{nearby}"')
            lines.append('')

        lines.append('### Rules')
        lines.append('')
        lines.append('- The logo must appear on **every slide** at the exact position above.')
        lines.append('- Do not stretch, crop, or recolor the logo.')
        lines.append('- Maintain clear space of at least the logo height around all edges.')

        return '\n'.join(lines)

    def _section_backgrounds(self) -> str:
        master_bg = self.backgrounds.get('master', {})
        bg_type = master_bg.get('type', 'inherited')

        lines = ['## 10. Background', '']
        lines.append(f'**Master background type**: `{bg_type}`')

        if bg_type == 'solid':
            lines.append(f'**Color**: `{master_bg.get("color", "N/A")}`')
        elif bg_type == 'picture':
            lines.append(f'**Image**: `images/{master_bg.get("image_filename", "background.png")}`')
            lines.append('Apply a dark overlay (brightness 0.3) for text readability.')
        elif bg_type == 'gradient':
            stops = master_bg.get('stops', [])
            if stops:
                stop_str = ', '.join(f'{s.get("color", "?")} at {s.get("position", 0):.0%}'
                                     for s in stops)
                lines.append(f'**Gradient**: {stop_str}')

        lines.append('')
        lines.append('### Rules')
        lines.append('')
        lines.append('- Content slides: use the master background as-is.')
        lines.append('- Title / cover slides: may use a layout-specific background')
        lines.append('  (see `layout_details` in the manifest).')
        lines.append('- Never use a background that clashes with the color tokens.')

        return '\n'.join(lines)

    def _section_motion(self) -> str:
        return """## 11. Motion & Transitions

```yaml
transitions:
  default: "slide"            # slide | fade | none
  duration: 400ms
  easing: "ease-out"
fragment_animations:
  default: "fade-up"          # fade-up | fade-in | zoom-in
  duration: 300ms
  stagger: 100ms              # delay between successive items
```

### Rules

- Use **one** transition type across the entire presentation.
- Fragment animations (`{.click}`) should use consistent direction.
- Avoid bounce, spin, or other playful animations in professional decks.
- Canvas animations follow their own timing (see `:::canvas` blocks)."""

    def _section_design_source(self) -> Optional[str]:
        if not self.design_source:
            return None

        lines = ['## 12. External Design Sources', '']

        if 'sources' in self.design_source:
            for src in self.design_source['sources']:
                lines.append(f'- **{src.get("type", "unknown").title()}**: `{src.get("url", "")}`')
                if src.get('description'):
                    lines.append(f'  {src["description"]}')
        else:
            lines.append(f'- **{self.design_source.get("type", "unknown").title()}**: `{self.design_source.get("url", "")}`')
            if self.design_source.get('description'):
                lines.append(f'  {self.design_source["description"]}')

        lines.append('')
        lines.append('> When the PPTX-extracted values conflict with the external')
        lines.append('> design guide, **the external guide takes precedence** for')
        lines.append('> color tokens and typography.  The PPTX values serve as')
        lines.append('> baseline defaults.')

        return '\n'.join(lines)

    def _section_usage_rules(self) -> str:
        return """## Design Checklist

Before finalizing any slide, verify:

- [ ] Colors use semantic tokens from §2 (no hardcoded hex)
- [ ] Typography follows the type scale from §3
- [ ] Icons follow the style rules from §4
- [ ] Content stays within the safe area from §5
- [ ] Shapes use the border-radius scale from §6
- [ ] Decorative elements match §7 exactly
- [ ] Header/footer regions are untouched per §8
- [ ] Logo appears at the correct position per §9
- [ ] Background matches the master template per §10
- [ ] Transitions are consistent per §11"""


def _build_design_source(args) -> Optional[dict[str, Any]]:
    """Build design_source dict from CLI args (--figma, --stitch)."""
    sources = []

    if getattr(args, 'figma', None):
        sources.append({
            'type': 'figma',
            'url': args.figma,
            'description': 'Figma design file — open to inspect colors, typography, spacing, and component styles.',
        })

    if getattr(args, 'stitch', None):
        sources.append({
            'type': 'stitch',
            'url': args.stitch,
            'description': 'Google Stitch design guide — reference for layout grid, color tokens, and asset exports.',
        })

    if not sources:
        return None
    if len(sources) == 1:
        return sources[0]
    return {'sources': sources}


def main():
    """Main entry point for CLI."""
    parser = argparse.ArgumentParser(
        description='Extract theme information from PPTX files and generate CSS overrides.',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s presentation.pptx -o ./theme
  %(prog)s template.pptx --list-masters
  %(prog)s template.pptx --master 1 -o ./theme --json-only
  %(prog)s template.pptx -o ./theme --design-md
  %(prog)s template.pptx -o ./theme --design-md --figma "https://figma.com/file/abc123"
  %(prog)s template.pptx -o ./theme --design-md --stitch "https://stitch.google.com/..."
        """
    )

    parser.add_argument('pptx_path', help='Path to the PPTX file')
    parser.add_argument('-o', '--output', default='./pptx-theme',
                       help='Output directory (default: ./pptx-theme)')
    parser.add_argument('--master', type=int, default=0,
                       help='Slide master index to use (default: 0)')
    parser.add_argument('--list-masters', action='store_true',
                       help='List all slide masters and exit')
    parser.add_argument('--list-layouts', action='store_true',
                       help='List all layouts in selected master and exit')
    parser.add_argument('--json-only', action='store_true',
                       help='Generate only JSON manifest, skip CSS')
    parser.add_argument('--css-file', default='theme-override.css',
                       help='CSS output filename (default: theme-override.css)')
    parser.add_argument('--design-md', action='store_true',
                       help='Generate design.md design-system document')
    parser.add_argument('--design-md-file', default='design.md',
                       help='Design guide output filename (default: design.md)')

    # External design-guide references
    design_group = parser.add_argument_group(
        'design sources',
        'Reference external design guides. URLs are stored in '
        'slide_master.design_source so agents can consult them '
        'when generating slides.'
    )
    design_group.add_argument(
        '--figma', metavar='URL',
        help='Figma file URL to reference as design guide')
    design_group.add_argument(
        '--stitch', metavar='URL',
        help='Google Stitch design guide URL to reference')

    args = parser.parse_args()

    try:
        extractor = ThemeExtractor(args.pptx_path, args.master)

        # Handle list options
        if args.list_masters:
            masters = extractor.list_masters()
            print(f"\nSlide Masters in {args.pptx_path}:")
            print("-" * 50)
            for m in masters:
                print(f"  [{m['index']}] {m['name']} ({m['layout_count']} layouts)")
            return 0

        if args.list_layouts:
            layouts = extractor.list_layouts()
            print(f"\nLayouts in Master {args.master}:")
            print("-" * 50)
            for l in layouts:
                print(f"  [{l['index']:2d}] {l['name']:<30} (bg: {l['background_type']})")
            return 0

        # Extract theme
        output_dir = Path(args.output)
        design_source = _build_design_source(args)

        print(f"Extracting theme from: {args.pptx_path}")
        print(f"Output directory: {output_dir}")
        if design_source:
            src_type = design_source.get('type', 'multiple')
            print(f"Design source: {src_type}")

        manifest = extractor.extract_all(output_dir, design_source=design_source)

        # Save JSON manifest
        manifest_path = output_dir / 'theme-manifest.json'
        with open(manifest_path, 'w') as f:
            json.dump(manifest, f, indent=2)
        print(f"Created: {manifest_path}")

        # Generate CSS unless json-only
        if not args.json_only:
            css_generator = CSSGenerator(manifest)
            css_content = css_generator.generate()

            css_path = output_dir / args.css_file
            with open(css_path, 'w') as f:
                f.write(css_content)
            print(f"Created: {css_path}")

        # Generate design.md design-system document
        if args.design_md:
            design_gen = DesignGuideGenerator(manifest)
            design_content = design_gen.generate()

            design_path = output_dir / args.design_md_file
            with open(design_path, 'w') as f:
                f.write(design_content)
            print(f"Created: {design_path}")

        # Summary
        print(f"\nTheme: {manifest['theme_name']}")
        print(f"Colors extracted: {len(manifest['colors'])}")
        print(f"Logos extracted: {len(manifest['logos'])}")
        print(f"Layouts: {manifest['layout_count']}")
        if design_source:
            print(f"Design source(s) linked: yes")
        if args.design_md:
            print(f"Design guide: {args.design_md_file}")

        return 0

    except FileNotFoundError as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1
    except IndexError as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        return 1


if __name__ == '__main__':
    sys.exit(main())
