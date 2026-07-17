#!/usr/bin/env python3
"""
Export a built reactive-presentation project to a screenshot-based .pptx
(one full-bleed image per slide — NOT an editable/native deck; for editable
PPTX use the aws-light-fcd skill).

Trust boundary: run this only on decks YOU built with this skill — the
project's HTML/JS executes in the headless browser during capture.

Renders each slide headlessly (Playwright/Chromium), captures a pixel-exact
screenshot with all fragments revealed and canvas step animations completed,
and assembles a PowerPoint (python-pptx) with speaker notes carried over
from the deck's presenter notes.

This is the reliable, scriptable counterpart to the in-browser
`ExportUtils.exportPPTX()` button (html2canvas): real browser rendering
means fonts, canvas drawings, shadows, and gradients survive intact.

Usage:
    python3 export_pptx.py <project-dir> [-o out.pptx] [--blocks a.html b.html]
                           [--width 1920] [--height 1080] [--scale 2]

    <project-dir>  Built output dir (contains index.html / NN-block.html + common/)

Dependencies:
    pip install 'playwright>=1.40' 'python-pptx>=1.0'   (+ `playwright install chromium`,
    or set PLAYWRIGHT_BROWSERS_PATH / CHROMIUM_PATH to an existing Chromium)
"""

import argparse
import functools
import http.server
import os
import re
import shutil
import sys
import tempfile
import threading
import urllib.parse
from pathlib import Path


def _die(msg: str, code: int = 1):
    print(f"ERROR: {msg}", file=sys.stderr)
    sys.exit(code)


def _import_deps():
    try:
        from playwright.sync_api import sync_playwright  # noqa: F401
    except ImportError:
        _die("playwright is not installed. Run: pip install playwright && playwright install chromium")
    try:
        import pptx  # noqa: F401
    except ImportError:
        _die("python-pptx is not installed. Run: pip install python-pptx")


class _QuietHandler(http.server.SimpleHTTPRequestHandler):
    def log_message(self, *args, **kwargs):
        pass


def _serve(root: Path):
    """Serve project dir on a free localhost port (file:// breaks fetch/CORS)."""
    handler = functools.partial(_QuietHandler, directory=str(root))

    # Bind port 0 directly on the server socket — no probe/rebind race.
    httpd = http.server.ThreadingHTTPServer(('127.0.0.1', 0), handler)
    port = httpd.server_address[1]
    thread = threading.Thread(target=httpd.serve_forever, daemon=True)
    thread.start()
    return httpd, f"http://127.0.0.1:{port}"


def _discover_blocks(project_dir: Path):
    """Pick HTML files to export: merged index.html if present, else block files."""
    index = project_dir / 'index.html'
    if index.exists() and 'name="generator" content="remarp"' in index.read_text(encoding='utf-8', errors='ignore'):
        return ['index.html']
    blocks = []
    for f in sorted(project_dir.glob('*.html')):
        if f.name in ('toc.html', 'index.html'):
            continue
        text = f.read_text(encoding='utf-8', errors='ignore')
        if 'name="generator" content="remarp"' in text:
            blocks.append(f.name)
    return blocks


# Runs inside the page: freeze motion, reveal fragments, finish canvas steps,
# hide navigation chrome that has no meaning on a static slide. The sidebar
# must also go: when visible it shrinks .slide-deck to calc(100vw - 220px),
# so a 1920px viewport would capture a ~1700px deck.
_PREPARE_JS = """
() => {
  const style = document.createElement('style');
  style.textContent = `
    *, *::before, *::after { transition: none !important; animation: none !important; }
    .progress-bar, .slide-counter, .nav-hint, .canvas-controls, .slide-sidebar { display: none !important; }
  `;
  document.head.appendChild(style);
  document.body.classList.remove('sidebar-visible');

  // Lazy iframes inside display:none slides never load before capture —
  // force eager loading so @type: iframe slides don't export blank. A load
  // flag is recorded so _MEDIA_READY_JS can wait even for cross-origin
  // frames (whose contentDocument is null and can't be inspected).
  document.querySelectorAll('iframe[loading="lazy"]').forEach(f => {
    f.dataset.exportLazy = '1';
    // A lazy frame that was already visible may have loaded before this runs
    // — the load event won't re-fire, so detect it now (same-origin only).
    try {
      const doc = f.contentDocument;
      if (doc && doc.readyState === 'complete' && doc.location.href !== 'about:blank') {
        f.dataset.exportLoaded = '1';
      }
    } catch (e) { /* cross-origin: leave to the load event */ }
    f.addEventListener('load', () => { f.dataset.exportLoaded = '1'; }, { once: true });
    f.loading = 'eager';
  });

  document.querySelectorAll('.fragment').forEach(el => el.classList.add('visible'));

  document.querySelectorAll('.slide').forEach(slide => {
    if (typeof slide.__canvasStep === 'function') {
      const parsed = parseInt(slide.dataset.canvasMaxStep || '60', 10);
      const max = Number.isFinite(parsed) ? parsed : 60;
      for (let s = 0; s < max; s++) {
        if (slide.__canvasStep('next') === false) break;
      }
    }
  });

  // Hand-authored canvases may register via deck.registerSlideAction()
  // instead of slide.__canvasStep — advance those to their final step too.
  try {
    if (typeof deck !== 'undefined' && deck && deck.slideActions) {
      Object.values(deck.slideActions).forEach(h => {
        if (!h || typeof h.down !== 'function') return;
        for (let s = 0; s < 60; s++) {
          if (h.down() === false) break;
        }
      });
    }
  } catch (e) { /* best effort */ }

  return document.querySelectorAll('.slide-deck .slide').length;
}
"""

_SHOW_SLIDE_JS = """
(idx) => {
  const slides = document.querySelectorAll('.slide-deck .slide');
  slides.forEach((el, i) => {
    el.classList.remove('entering', 'leaving');
    if (i === idx) {
      el.classList.add('active');
      el.style.display = 'flex';
      el.style.opacity = '1';
      el.style.visibility = 'visible';
    } else {
      el.classList.remove('active');
      el.style.display = 'none';
    }
  });
  // We bypass SlideFramework.showSlide(), so replicate its per-slide chrome
  // updates: pagination number, and footer/logo visibility + dark-logo swap
  // (updateFooterVisibility hides the framework footer/logo on slides with
  // an <img> and swaps logoDarkSrc on dark slides — without this, slide 0's
  // state is baked into every capture).
  const num = document.querySelector('.slide-number');
  if (num) num.textContent = (idx + 1) + ' / ' + slides.length;
  const el = slides[idx];
  try {
    if (typeof deck !== 'undefined' && deck &&
        typeof deck.updateFooterVisibility === 'function' && el) {
      deck.updateFooterVisibility(el);
    }
  } catch (e) { /* best effort */ }
  const h = el ? el.querySelector('h1, h2, h3') : null;
  return h ? h.textContent.trim() : '';
}
"""

# Media inside the ACTIVE slide (images, iframes) must finish loading before
# capture; cross-origin iframes can't be inspected and are assumed ready.
_MEDIA_READY_JS = """
() => {
  const el = document.querySelector('.slide-deck .slide.active');
  if (!el) return true;
  const imgs = Array.from(el.querySelectorAll('img'));
  const iframes = Array.from(el.querySelectorAll('iframe'));
  return imgs.every(i => i.complete) && iframes.every(f => {
    // Frames we forced from lazy→eager must fire their load event first —
    // this also covers cross-origin frames, which expose no readyState.
    if (f.dataset.exportLazy === '1' && f.dataset.exportLoaded !== '1') return false;
    try {
      const doc = f.contentDocument;
      if (!doc) return true;  // cross-origin, already loaded before prepare
      return doc.readyState === 'complete';
    } catch (e) { return true; }
  });
}
"""

# presenterNotes is a top-level `const` in the built HTML; code evaluated in
# the page's global scope (like the DevTools console) can still read it.
_NOTES_JS = """
() => {
  try { return typeof presenterNotes !== 'undefined' ? presenterNotes : {}; }
  catch (e) { return {}; }
}
"""


def _notes_to_text(text: str) -> str:
    """Presenter notes → PPTX notes text.

    presenterNotes carries the RAW :::notes markdown (plus [timing]/cue
    markers) — it is not HTML, so no tag stripping: a generic <[^>]+> pass
    would eat real content like `List<T>`. Only normalize whitespace.
    """
    text = str(text).replace('\r\n', '\n').replace('\r', '\n')
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def export(project_dir: Path, out_path: Path, blocks, width: int, height: int, scale: int):
    from playwright.sync_api import sync_playwright
    from pptx import Presentation
    from pptx.util import Inches

    httpd, base_url = _serve(project_dir)
    tmpdir = Path(tempfile.mkdtemp(prefix='remarp-pptx-'))

    prs = Presentation()
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide_size_set = False

    total = 0
    try:
        with sync_playwright() as p:
            launch_kwargs = {}
            chromium_path = os.environ.get('CHROMIUM_PATH')
            if chromium_path:
                launch_kwargs['executable_path'] = chromium_path
            browser = p.chromium.launch(**launch_kwargs)
            page = browser.new_page(
                viewport={'width': width, 'height': height},
                device_scale_factor=scale,
            )

            for block in blocks:
                page.goto(f"{base_url}/{urllib.parse.quote(block)}", wait_until='networkidle')
                page.evaluate("() => document.fonts.ready.then(() => true)")
                page.wait_for_timeout(600)  # let canvas setup scripts settle

                slide_count = page.evaluate(_PREPARE_JS)
                if not slide_count:
                    print(f"  (warn) no slides found in {block} — skipping")
                    continue
                notes = page.evaluate(_NOTES_JS) or {}
                if not notes:
                    print(f"  (warn) no presenter notes found in {block}")
                page.wait_for_timeout(300)  # canvas redraw after step advance

                deck = page.locator('.slide-deck')

                # PPTX slide size must match what is actually captured: the
                # deck enforces ITS OWN aspect ratio via CSS max-width/height,
                # so a 4:3 viewport still captures a 16:9 deck (and vice
                # versa). Deriving the size from the viewport would stretch
                # the image — derive it from the deck's real bounding box.
                if not slide_size_set:
                    box = deck.bounding_box()
                    if box and box['width'] and box['height']:
                        deck_ratio = box['width'] / box['height']
                        if abs(deck_ratio - width / height) > 0.01:
                            print(f"  (warn) viewport {width}x{height} does not match the "
                                  f"deck's own aspect ratio ({deck_ratio:.3f}) — "
                                  f"slide size follows the deck")
                        prs.slide_width = Inches(round(7.5 * deck_ratio, 3))
                    else:
                        prs.slide_width = Inches(round(7.5 * width / height, 3))
                    slide_size_set = True

                for i in range(slide_count):
                    title = page.evaluate(_SHOW_SLIDE_JS, i)
                    # Media inside the now-visible slide (images, formerly
                    # lazy iframes) must finish loading before capture.
                    try:
                        page.wait_for_function(_MEDIA_READY_JS, timeout=5000)
                    except Exception:
                        print(f"  (warn) media still loading on slide {i + 1} of {block} — capturing anyway")
                    page.wait_for_timeout(80)

                    img = tmpdir / f"{Path(block).stem}-{i + 1:03d}.png"
                    deck.screenshot(path=str(img))

                    pptx_slide = prs.slides.add_slide(blank)
                    pptx_slide.shapes.add_picture(
                        str(img), 0, 0, prs.slide_width, prs.slide_height)

                    note = notes.get(str(i + 1)) or notes.get(i + 1) or ''
                    note_text = _notes_to_text(note)
                    if title and note_text:
                        note_text = f"{title}\n\n{note_text}"
                    elif title:
                        note_text = title
                    if note_text:
                        pptx_slide.notes_slide.notes_text_frame.text = note_text

                    total += 1
                    print(f"  [{total}] {block} slide {i + 1}/{slide_count}"
                          + (f" — {title}" if title else ""))

            browser.close()
    finally:
        httpd.shutdown()
        shutil.rmtree(tmpdir, ignore_errors=True)

    if total == 0:
        _die("no slides were captured — nothing to export "
             "(are the block files built remarp HTML?)")

    prs.save(str(out_path))
    print(f"\nOK: {total} slides → {out_path}")


def main():
    ap = argparse.ArgumentParser(description='Export built reactive-presentation to .pptx')
    ap.add_argument('project_dir', help='Built output dir (contains *.html + common/)')
    ap.add_argument('-o', '--output', help='Output .pptx path (default: <dir-name>.pptx in project dir)')
    ap.add_argument('--blocks', nargs='+', help='Specific block HTML files to export (relative to project dir)')
    ap.add_argument('--width', type=int, default=1920, help='Capture viewport width (default 1920)')
    ap.add_argument('--height', type=int, default=1080, help='Capture viewport height (default 1080)')
    ap.add_argument('--scale', type=int, default=2, choices=(1, 2, 3),
                    help='Device scale factor for crisp captures (default 2)')
    args = ap.parse_args()

    _import_deps()

    project_dir = Path(args.project_dir).resolve()
    if not project_dir.is_dir():
        _die(f"not a directory: {project_dir}")
    if args.width <= 0 or args.height <= 0:
        _die("--width/--height must be positive")
    if not (project_dir / 'common').is_dir():
        print("  (warn) no common/ directory inside the project dir — built decks "
              "are self-contained (build copies common/ into the output dir); "
              "point this script at the BUILT deck directory or assets will 404")

    blocks = args.blocks or _discover_blocks(project_dir)
    if not blocks:
        _die(f"no remarp HTML files found in {project_dir} (build first: remarp_to_slides.py build)")
    for b in blocks:
        if not (project_dir / b).exists():
            _die(f"block not found: {b}")

    out_path = Path(args.output) if args.output else project_dir / f"{project_dir.name}.pptx"

    print(f"Exporting {len(blocks)} file(s) from {project_dir} → {out_path}")
    export(project_dir, out_path, blocks, args.width, args.height, args.scale)


if __name__ == '__main__':
    main()
